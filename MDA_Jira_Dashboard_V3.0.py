
# -*- coding: utf-8 -*-
# MDA_Jira_Dashboard_V3.0.py  —  SELF-CONTAINED (no external .py dependencies)
# All logic from Combined_Automation, MDA_Manuj_Weekly_Report, and slide2_builder
# is bundled inline. Only standard pip packages are required.

import os, re, io, sys, types, traceback, time, copy, inspect
from datetime import datetime, timedelta
from pathlib import Path
from contextlib import redirect_stdout
from html import escape
from dataclasses import dataclass
from typing import Optional, List, Tuple

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
from playwright.sync_api import sync_playwright, TimeoutError as PWTimeoutError

if sys.platform.startswith('win'):
    try:
        import asyncio
        asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
    except Exception:
        pass

BASE_DIR = Path(__file__).parent
ARTIFACTS = (BASE_DIR / 'artifacts'); ARTIFACTS.mkdir(exist_ok=True)
DOWNLOADS_DIR = Path.home() / 'Downloads'; DOWNLOADS_DIR.mkdir(exist_ok=True)

# ═══════════════════════════════════════════════════════════════════════════════
# EMBEDDED: Combined_Automation_Filtered_iterPrompt_v3.py
# ═══════════════════════════════════════════════════════════════════════════════

JIRA_BASE_URL = 'https://jira.charter.com'
JIRA_USER = os.getenv('JIRA_USER', '')
JIRA_PASS = os.getenv('JIRA_PASS', '')
CA_HEADLESS = os.getenv('HEADLESS', 'false').lower() == 'true'
CA_TIMEOUT_MS = int(os.getenv('TIMEOUT_MS', '45000'))

ALLOWED_ASSIGNEES = ['JIRAUSER776479', 'JIRAUSER726004', 'JIRAUSER775575', 'JIRAUSER642656',
                     'JIRAUSER640395', 'JIRAUSER747252', 'JIRAUSER776120']
ALLOWED_REPORTERS = ['Dorairaj, Susintha', 'Rout, Swaha', 'Shinde, Pratiksha', 'Patil, Kalawati',
                     'Khatoon, Asma', 'Kannydhary, Jaswanth', 'Kappagantula, Surya Prakash']
ALLOWED_REPORTERS_ENV = os.getenv('ALLOWED_REPORTERS_CSV', '')
if ALLOWED_REPORTERS_ENV:
    ALLOWED_REPORTERS = [s.strip() for s in ALLOWED_REPORTERS_ENV.split(',') if s.strip()]
NAME_MAPPING = {
    'JIRAUSER776479': 'Surya', 'JIRAUSER726004': 'Susintha', 'JIRAUSER775575': 'Asma',
    'JIRAUSER642656': 'Prathiksha', 'JIRAUSER640395': 'Swaha', 'JIRAUSER747252': 'Kalawathi',
    'JIRAUSER776120': 'Jashwanth'
}

def _ca_wait_idle(page, timeout=CA_TIMEOUT_MS):
    try: page.wait_for_load_state('networkidle', timeout=timeout)
    except Exception: pass

def _ca_smart_wait(page, selector, timeout=10000):
    try: page.wait_for_selector(selector, state='visible', timeout=timeout); return True
    except Exception: return False

def _ca_wait_page_ready(page, timeout=30000):
    try: page.wait_for_function("document.readyState === 'complete'", timeout=timeout)
    except Exception: pass

def _ca_click_element(page, *selectors, timeout=5000):
    for sel in selectors:
        try: page.wait_for_selector(sel, timeout=timeout); page.locator(sel).first.click(); return True
        except Exception: continue
    return False

def _ca_read_defect_excel(excel_file: Path) -> pd.DataFrame:
    try: return pd.read_excel(excel_file, engine='openpyxl')
    except Exception: pass
    try: return pd.read_excel(excel_file, engine='xlrd')
    except Exception: pass
    tables = pd.read_html(str(excel_file), flavor='lxml')
    return tables[1] if len(tables) > 1 else tables[0]

def filter_qmetry_data(excel_file: Path) -> Path:
    df = pd.read_excel(excel_file)
    original_count = len(df)
    if 'Assignee' in df.columns:
        df_filtered = df[df['Assignee'].isin(ALLOWED_ASSIGNEES)]
        df_filtered['Assignee'] = df_filtered['Assignee'].map(NAME_MAPPING).fillna(df_filtered['Assignee'])
        print(f'Filtered QMetry: {original_count} -> {len(df_filtered)} test cases')
        filtered_path = excel_file.parent / f'{excel_file.stem}_filtered{excel_file.suffix}'
        df_filtered.to_excel(filtered_path, index=False)
        return filtered_path
    else:
        print("Warning: 'Assignee' column not found in QMetry export")
        return excel_file

def filter_defect_data(excel_file: Path) -> Path:
    df = _ca_read_defect_excel(excel_file)
    original_count = len(df)
    if 'Reporter' in df.columns:
        df_filtered = df[df['Reporter'].isin(ALLOWED_REPORTERS)]
        print(f'Filtered Defects: {original_count} -> {len(df_filtered)} defects')
        filtered_path = excel_file.parent / f'{excel_file.stem}_filtered.xlsx'
        df_filtered.to_excel(filtered_path, index=False)
        return filtered_path
    else:
        print("Warning: 'Reporter' column not found in Defect export")
        return excel_file

def create_qmetry_pivot(excel_file: Path, integration_name: str) -> Path:
    df = pd.read_excel(excel_file)
    if 'Assignee' in df.columns:
        df['Assignee'] = df['Assignee'].map(NAME_MAPPING).fillna(df['Assignee'])
    if 'Labels' not in df.columns:
        raise ValueError("QMetry export is missing 'Labels' column")
    grouped = df.groupby('Labels').agg(
        Count=('Issue Key', 'count'),
        Assignee=('Assignee', lambda x: ', '.join(pd.Series(x).dropna().unique()))
    ).reset_index()
    pivot = pd.DataFrame({
        'Release': [f'{integration_name} - {label}' for label in grouped['Labels']],
        'Project Name': 'T-Mobile MVNO (MVP)',
        'Offshore TC execution By': grouped['Assignee'].values,
        'Test case creation Count': grouped['Count'].values,
        'Reviewer name': 'Yerra,Bijay',
        'Comments incorporated': 'Closed',
        'Review with Product': 'Completed',
        'Comments': 'No Comments'
    })
    pivot_path = ARTIFACTS / f'pivot_labels_testcases_{excel_file.stem}.xlsx'
    with pd.ExcelWriter(pivot_path, engine='openpyxl') as writer:
        pivot.to_excel(writer, index=False, sheet_name='Labels Pivot')
        ws = writer.sheets['Labels Pivot']
        for column in ws.columns:
            ws.column_dimensions[column[0].column_letter].width = max(len(str(cell.value)) for cell in column) + 2
    print(f'QMetry Pivot saved: {pivot_path}')
    return pivot_path

def create_defect_pivot(excel_file: Path, integration_name: str) -> Path:
    df = _ca_read_defect_excel(excel_file)
    if 'Status' not in df.columns or 'Priority' not in df.columns or 'Key' not in df.columns:
        raise ValueError('Defect export missing one of required columns: Status / Priority / Key')
    pivot = pd.pivot_table(df, values='Key', index='Status', columns='Priority', aggfunc='count', fill_value=0)
    pivot.reset_index(inplace=True)
    pivot.insert(0, 'Release', integration_name)
    pivot_path = ARTIFACTS / f'pivot_{Path(excel_file).stem}.xlsx'
    with pd.ExcelWriter(pivot_path, engine='openpyxl') as writer:
        pivot.to_excel(writer, index=False, sheet_name='Defects Pivot')
        ws = writer.sheets['Defects Pivot']
        if len(pivot) > 1:
            ws.merge_cells(start_row=2, start_column=1, end_row=len(pivot)+1, end_column=1)
        for column in ws.columns:
            ws.column_dimensions[column[0].column_letter].width = max(len(str(cell.value)) for cell in column) + 2
    print(f'Defect Pivot saved: {pivot_path}')
    return pivot_path

def run_qmetry_automation(browser, iteration: str, integration_name: str):
    context = browser.new_context(accept_downloads=True)
    context.tracing.start(screenshots=True, snapshots=True)
    page = context.new_page()
    try:
        print('\n=== QMetry Test Case Export ===')
        print('1. Navigating to Jira...')
        page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/?projectId=127207', timeout=120000, wait_until='domcontentloaded')
        _ca_wait_page_ready(page); _ca_wait_idle(page, timeout=30000)
        if page.locator("input[type='password']").count() > 0:
            if JIRA_USER and JIRA_PASS:
                page.fill("input[name='username'], input[type='text']", JIRA_USER)
                page.fill("input[type='password']", JIRA_PASS)
                page.click("button[type='submit']"); _ca_wait_idle(page)
            else:
                input('Login manually, then press Enter...')
        print('2. Navigating to QMetry...')
        _ca_smart_wait(page, "//a[contains(text(), 'More')]", timeout=15000)
        page.locator("//a[contains(text(), 'More')]").first.click()
        _ca_smart_wait(page, "//a[contains(text(), 'QMetry')]", timeout=5000)
        page.locator("//a[contains(text(), 'QMetry')]").first.hover()
        _ca_smart_wait(page, "//a[contains(text(), 'Test Management')]", timeout=5000)
        page.locator("//a[contains(text(), 'Test Management')]").first.click()
        _ca_wait_idle(page)
        print('3. Selecting Software Project...')
        _ca_smart_wait(page, "//span[contains(text(), 'Software Project')]", timeout=10000)
        page.locator("//span[contains(text(), 'Software Project')]").first.click()
        page.keyboard.type('MDA-Integration'); page.keyboard.press('Enter')
        _ca_wait_idle(page, timeout=10000)
        print('4. Searching for folder...')
        _ca_smart_wait(page, 'div.glyph-search', timeout=10000)
        page.locator('div.glyph-search').first.click()
        folder_name = f'INTEGRATION_PROGRESSION_{iteration}'
        page.keyboard.type(folder_name); page.keyboard.press('Enter')
        _ca_wait_idle(page, timeout=15000)
        print('5. Clicking xpath element...')
        xpath5 = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div/div[1]/div[2]/div[1]/div/button/span/span/span"
        _ca_smart_wait(page, xpath5, timeout=10000); page.locator(xpath5).first.click()
        print('6. Checking checkbox...')
        checkbox_xpath = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div/div[1]/div[2]/div[2]/div/div/ul/li[3]/label/span[1]/span/span"
        try:
            page.wait_for_selector(checkbox_xpath, timeout=5000)
            checkbox = page.locator(checkbox_xpath).first
            if not checkbox.is_checked(): checkbox.click()
        except Exception:
            page.locator(checkbox_xpath).first.click()
        print('6b. Waiting for page to fully load after checkbox...')
        _ca_wait_page_ready(page, timeout=60000)
        _ca_wait_idle(page, timeout=30000)
        page.wait_for_timeout(3000)
        print('7. Scrolling to find CHPROJECT-18431...')
        target_sel = "//span[contains(text(), 'CHPROJECT-18431')]"
        # Try multiple scrollable containers
        scroll_selectors = [
            "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div/ul/div/div",
            "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div/ul",
            "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div",
            "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]",
            "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav",
        ]
        found = _ca_smart_wait(page, target_sel, timeout=5000)
        if not found:
            for sc_sel in scroll_selectors:
                try:
                    sc = page.locator(sc_sel).first
                    if sc.count() == 0: continue
                    print(f'   Trying scroll container: {sc_sel[-60:]}')
                    for i in range(20):
                        sc.evaluate(f'el => el.scrollTop += 300')
                        page.wait_for_timeout(500)
                        if _ca_smart_wait(page, target_sel, timeout=1500):
                            found = True; print(f'   Found after scroll #{i+1}'); break
                    if found: break
                    # Reset and try next container
                    sc.evaluate('el => el.scrollTop = 0')
                except Exception as e:
                    print(f'   [WARN] Scroll container failed: {e}')
                    continue
        if not found:
            # Last resort: use keyboard to scroll the page
            print('   Trying keyboard scroll fallback...')
            page.keyboard.press('End')
            page.wait_for_timeout(2000)
            found = _ca_smart_wait(page, target_sel, timeout=5000)
        if not found:
            page.screenshot(path=str(ARTIFACTS / 'qmetry_scroll_debug.png'))
            raise RuntimeError('CHPROJECT-18431 not found after exhaustive scrolling')
        page.locator(target_sel).first.click(button='right')
        print('8. Clicking Export to Excel...')
        _ca_smart_wait(page, "//span[contains(text(), 'Export to Excel')]", timeout=5000)
        page.locator("//span[contains(text(), 'Export to Excel')]").first.click()
        print('9. Downloading file...')
        notif_xpath = "//*[@id='root_4.14.1.1']/div/div/div/div[1]/div[2]/div/div[3]/div/span/span"
        if _ca_smart_wait(page, notif_xpath, timeout=45000):
            try: page.locator(notif_xpath).first.click()
            except Exception: pass
        dl_btn_selectors = [
            "//*[@id='root_4.14.1.1']/div[2]/div[2]/div/div/div[2]/div[1]/div[1]/div[2]/div[1]/div[1]/div/div/button/span/span[1]",
            "//button[contains(@class, 'download') or contains(@title, 'Download') or contains(@aria-label, 'Download')]",
            "//button[.//span[contains(text(), 'Download') or contains(text(), 'download')]]",
            "//a[contains(@class, 'download') or contains(@title, 'Download')]",
            "//button[contains(@class, 'btn')][.//span[contains(@class, 'icon-download') or contains(@class, 'glyphicon-download')]]",
        ]
        element = None
        for sel in dl_btn_selectors:
            if _ca_smart_wait(page, sel, timeout=15000):
                element = page.locator(sel).first; print(f'   Download button found via: {sel[:70]}'); break
        if element is None:
            page.screenshot(path=str(ARTIFACTS / 'qmetry_download_debug.png'))
            for sel in ["button:has-text('Download')", "a:has-text('Download')", "[role='dialog'] button"]:
                try: page.wait_for_selector(sel, timeout=5000); element = page.locator(sel).first; break
                except Exception: continue
        if element is None:
            raise RuntimeError('Could not find download button after export')
        # Wait for button to become enabled (export file preparation)
        print('   Waiting for download button to become enabled...')
        for _wait_i in range(60):
            is_disabled = element.evaluate('el => el.disabled || el.getAttribute("aria-disabled") === "true"')
            if not is_disabled:
                print(f'   Button enabled after {_wait_i + 1}s')
                break
            page.wait_for_timeout(1000)
        else:
            print('   [WARN] Button still disabled after 60s, attempting click anyway')
        element.hover()
        with page.expect_download(timeout=60000) as download_info:
            element.click()
        download = download_info.value
        download_path = ARTIFACTS / download.suggested_filename
        download.save_as(download_path)
        print(f'Downloaded: {download_path}')
        try: context.tracing.stop(path=ARTIFACTS / 'qmetry_trace.zip')
        except Exception: pass
        context.close()
        print('10. Filtering and creating pivot table...')
        filtered_path = filter_qmetry_data(download_path)
        pivot_path = create_qmetry_pivot(filtered_path, integration_name)
        print('QMetry automation completed!\n')
        return pivot_path
    except Exception as e:
        print(f'\n\u2716 QMetry Error: {e}')
        traceback.print_exc()
        try: page.screenshot(path=str(ARTIFACTS / 'qmetry_error.png'))
        except Exception: pass
        try: context.tracing.stop(path=ARTIFACTS / 'qmetry_trace.zip')
        except Exception: pass
        try: context.close()
        except Exception: pass
        return None

def run_defect_automation(browser, iteration: str, integration_name: str):
    context = browser.new_context(accept_downloads=True)
    context.tracing.start(screenshots=True, snapshots=True)
    page = context.new_page()
    page.set_default_timeout(60000)
    try:
        print('\n=== Jira Defect Export ===')
        print('1. Navigating to Jira...')
        page.goto(JIRA_BASE_URL, timeout=120000, wait_until='domcontentloaded')
        _ca_wait_idle(page)
        if page.locator("input[type='password']").count() > 0:
            if JIRA_USER and JIRA_PASS:
                page.fill("input[name='username'], input[type='text']", JIRA_USER)
                page.fill("input[type='password']", JIRA_PASS)
                page.click("button[type='submit']"); _ca_wait_idle(page)
            else:
                input('Login manually, then press Enter...')
        print('2. Navigating to Issues -> Search for issues...')
        _ca_smart_wait(page, "//a[contains(text(), 'Issues')]", timeout=15000)
        _ca_click_element(page, "//a[contains(text(), 'Issues')]", "a[href*='issues']", timeout=10000)
        _ca_smart_wait(page, "//a[contains(text(), 'Search for issues')]", timeout=10000)
        _ca_click_element(page, "//a[contains(text(), 'Search for issues')]", "a[href*='IssueNavigator']", timeout=10000)
        _ca_wait_idle(page)
        print('3. Clicking advanced search...')
        for sel in ["//*[@id='advanced-search']", "//a[contains(text(), 'Advanced')]",
                    "//a[contains(@class, 'switcher-item') and contains(text(), 'Advanced')]"]:
            try: page.wait_for_selector(sel, timeout=5000); page.locator(sel).first.click(); break
            except Exception: continue
        _ca_smart_wait(page, "textarea#advanced-search, textarea[name='jqlQuery'], #jqltext", timeout=8000)
        print('4. Entering search query...')
        search_query = f'project = "WPATEST-Integration Test Team" AND issuetype = Defect AND labels = INTG_UAT_PROGRESSION_{iteration} and "Service Type." = T-Mobile'
        for sel in ['textarea#advanced-search', 'textarea[name=\'jqlQuery\']', '#jqltext', '#advanced-search']:
            try:
                page.wait_for_selector(sel, timeout=5000)
                page.locator(sel).first.fill(search_query); page.keyboard.press('Enter')
                print(f'   JQL entered via: {sel}'); break
            except Exception: continue
        _ca_wait_idle(page)
        _ca_smart_wait(page, '.issuetable, .issue-table, #issuetable, .no-results-message, .results-count', timeout=15000)
        print('5. Fetching defects via REST API...')
        rest_url = f'{JIRA_BASE_URL}/rest/api/2/search'
        all_issues = []; start_at = 0
        while True:
            api_result = page.evaluate("""async ([url, jql, startAt]) => {
                const params = new URLSearchParams({jql: jql, startAt: startAt, maxResults: 100,
                    fields: 'key,summary,status,priority,reporter,assignee,issuetype,created,updated,labels'});
                const resp = await fetch(url + '?' + params.toString(), {
                    credentials: 'include', headers: {'Accept': 'application/json', 'Content-Type': 'application/json'}});
                if (!resp.ok) return {status: resp.status, body: null};
                return {status: resp.status, body: await resp.json()};
            }""", [rest_url, search_query, str(start_at)])
            if api_result['status'] != 200 or not api_result['body']:
                raise RuntimeError(f"REST API returned HTTP {api_result['status']}")
            data = api_result['body']; issues = data.get('issues', [])
            all_issues.extend(issues); total = data.get('total', 0)
            print(f'   Fetched {len(all_issues)}/{total} issues...')
            if len(all_issues) >= total or not issues: break
            start_at += len(issues)
        rows = []
        for issue in all_issues:
            fields = issue.get('fields', {})
            rows.append({
                'Key': issue.get('key', ''), 'Summary': fields.get('summary', ''),
                'Status': (fields.get('status') or {}).get('name', ''),
                'Priority': (fields.get('priority') or {}).get('name', ''),
                'Reporter': (fields.get('reporter') or {}).get('displayName', ''),
                'Assignee': (fields.get('assignee') or {}).get('displayName', ''),
                'Issue Type': (fields.get('issuetype') or {}).get('name', ''),
                'Created': fields.get('created', ''), 'Updated': fields.get('updated', ''),
            })
        df = pd.DataFrame(rows)
        download_path = ARTIFACTS / f'Jira_defects_{iteration}.xlsx'
        df.to_excel(download_path, index=False)
        print(f'   Saved {len(df)} defects to: {download_path}')
        context.tracing.stop(path=ARTIFACTS / 'trace.zip'); context.close()
        print('6. Filtering and creating pivot table...')
        filtered_path = filter_defect_data(download_path)
        pivot_path = create_defect_pivot(filtered_path, integration_name)
        print('Defect automation completed!\n')
        return pivot_path
    except Exception as e:
        print(f'\n\u2716 Defect Error: {e}')
        traceback.print_exc()
        try: page.screenshot(path=str(ARTIFACTS / 'defect_error.png'))
        except Exception: pass
        try: context.tracing.stop(path=ARTIFACTS / 'trace.zip')
        except Exception: pass
        try: context.close()
        except Exception: pass
        return None

# ═══════════════════════════════════════════════════════════════════════════════
# EMBEDDED: MDA_Manuj_Weekly_Reportv1.3.py
# ═══════════════════════════════════════════════════════════════════════════════

NAVY = RGBColor(11, 29, 57)
WHITE = RGBColor(255, 255, 255)
FONT_NAME = 'Segoe UI'

def _mmu_log(msg: str):
    print(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}", flush=True)

def week_monday_friday(today=None):
    today = today or datetime.now()
    mon = today - timedelta(days=today.weekday())
    fri = mon + timedelta(days=4)
    return mon.strftime('%d %b %Y'), fri.strftime('%d %b %Y')

@dataclass
class LayoutPolicy:
    top_in: float = 1.05
    side_in: float = 0.90
    footer_in: float = 2.00
    title_gap_in: float = 0.45
    rows_per_part: int = 10
    header_pt: int = 10
    body_pt: int = 8

def chunk_rows(df: pd.DataFrame, size: int) -> List[pd.DataFrame]:
    if df is None or df.empty:
        return [pd.DataFrame({'No Data': ['—']})]
    n = len(df)
    if n <= size:
        return [df.reset_index(drop=True)]
    q, r = divmod(n, size)
    sizes: List[int] = []
    if r == 1 and q >= 1:
        if q == 1:
            sizes = [size - 1, 2]
        else:
            sizes = [size] * (q - 1) + [size - 1, 2]
    else:
        sizes = [size] * q + ([r] if r else [])
    parts: List[pd.DataFrame] = []
    i = 0
    for s in sizes:
        parts.append(df.iloc[i:i+s].reset_index(drop=True))
        i += s
    return parts

_NARROW_COL_TOKENS = {'COUNT', 'TOTAL', 'TOTALS', 'COUNT(S)', 'SUM', 'GRAND TOTAL',
                      'BLOCKED', 'P0', 'P1', 'P2', 'P3', 'P4', 'NO RUN', 'NOT APPLICABLE',
                      'BLOCKER/EMERGENCY', 'CRITICAL/HIGH', 'MAJOR/MEDIUM',
                      'MINOR/LOW', 'TRIVIAL', 'BLOCKER', 'CRITICAL', 'MAJOR',
                      'MINOR', 'EMERGENCY', 'HIGH', 'MEDIUM', 'LOW',
                      'PASS', 'FAIL', 'NOT EXECUTED', 'IN PROGRESS',
                      'WORK IN PROGRESS', 'CONDITIONAL PASS', 'DEFERRED'}

def _is_narrow_col(col_name_upper: str) -> bool:
    if col_name_upper in _NARROW_COL_TOKENS:
        return True
    for token in _NARROW_COL_TOKENS:
        if token in col_name_upper:
            return True
    return False

def content_weights(df: pd.DataFrame) -> List[float]:
    if df is None or df.empty:
        return [1.0]
    lens = []
    for c in df.columns:
        vals = [str(c)] + [str(v) for v in df[c].astype(str).values]
        lens.append(max(len(v) for v in vals))
    MIN_NARROW = 7; MAX_COL_LEN = 28
    if len(lens) > 1:
        for i in range(len(lens)):
            lens[i] = min(lens[i], MAX_COL_LEN)
            col_name = str(df.columns[i]).strip().upper()
            if _is_narrow_col(col_name):
                lens[i] = MIN_NARROW  # cap narrow cols to fixed width
        lens[-1] = max(lens[-1], MIN_NARROW)
    tot = sum(lens) or 1
    return [L / tot for L in lens]

def merge_release_column(tbl, part_df: pd.DataFrame, release_col_index: int = 0, body_font_pt: int = 8):
    """Merge consecutive identical Release cells using raw XML (rowSpan / vMerge)
    so that python-pptx never concatenates text from other columns."""
    if part_df is None or part_df.empty:
        return
    values = list(part_df.iloc[:, release_col_index].astype(str).fillna('').values)
    start = 0
    while start < len(values):
        end = start
        while end + 1 < len(values) and values[end + 1] == values[start] and values[start] != '':
            end += 1
        if end > start:
            span = end - start + 1
            # Set rowSpan on the first cell (table row = data row + 1 for header)
            first_tr = tbl._tbl.tr_lst[start + 1]
            first_tc = first_tr.findall(qn('a:tc'))[release_col_index]
            first_tc.set('rowSpan', str(span))
            # Set vMerge on subsequent cells
            for offset in range(1, span):
                tr = tbl._tbl.tr_lst[start + 1 + offset]
                tc = tr.findall(qn('a:tc'))[release_col_index]
                tc.set('vMerge', '1')
                # Clear text in merged-away cells
                txBody = tc.find(qn('a:txBody'))
                if txBody is not None:
                    for p in txBody.findall(qn('a:p')):
                        for r in p.findall(qn('a:r')):
                            t_el = r.find(qn('a:t'))
                            if t_el is not None:
                                t_el.text = ''
        start = end + 1
        start = end + 1

def add_brand_footer_logo(s, prs, logo_path, width_in: float = 1.5):
    x = prs.slide_width - Inches(0.90 + width_in)
    y = prs.slide_height - Inches(0.55 + 0.50)
    if logo_path and Path(logo_path).exists():
        try:
            s.shapes.add_picture(str(logo_path), x, y, width=Inches(width_in))
            return
        except Exception:
            pass
    box = s.shapes.add_textbox(x, y + Inches(0.10), Inches(width_in), Inches(0.5))
    tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = 'Charter'; p.font.bold = True; p.font.size = Pt(22); p.font.name = FONT_NAME; p.font.color.rgb = WHITE

def add_title_slide(prs, logo_path):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = bg.fill; fill.solid(); fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.9), Inches(1.9), prs.slide_width - Inches(1.8), Inches(1.6))
    tf = box.text_frame; tf.clear()
    p0 = tf.paragraphs[0]; p0.text = 'MDA QA(T-Mob & VZ)'; p0.alignment = PP_ALIGN.LEFT
    p0.font.name = FONT_NAME; p0.font.size = Pt(44); p0.font.bold = True; p0.font.color.rgb = WHITE
    p1 = tf.add_paragraph(); p1.text = 'TOSCA Weekly Status Report'; p1.alignment = PP_ALIGN.LEFT
    p1.font.name = FONT_NAME; p1.font.size = Pt(32); p1.font.bold = True; p1.font.color.rgb = WHITE
    mon, fri = week_monday_friday()
    b2 = s.shapes.add_textbox(Inches(0.9), Inches(3.7), prs.slide_width - Inches(1.8), Inches(0.6))
    tf2 = b2.text_frame; tf2.clear(); p = tf2.paragraphs[0]
    p.text = f'– {mon} to {fri}'; p.alignment = PP_ALIGN.LEFT
    p.font.name = FONT_NAME; p.font.size = Pt(20); p.font.color.rgb = WHITE
    add_brand_footer_logo(s, prs, logo_path)

def _set_cell_autofit(cell):
    """Set XML-level shrink-text-on-overflow so PPT never expands the row."""
    txBody = cell._tc.find(qn('a:txBody'))
    if txBody is None:
        return
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
        txBody.insert(0, bodyPr)
    # Remove any existing autofit / spAutoFit
    for child in list(bodyPr):
        if child.tag in (qn('a:normAutofit'), qn('a:spAutoFit'), qn('a:noAutofit')):
            bodyPr.remove(child)
    etree.SubElement(bodyPr, qn('a:normAutofit')).set('fontScale', '50000')  # allow shrink to 50%
    # Minimal margins for tight rows
    bodyPr.set('lIns', '18288')   # ~0.02in left
    bodyPr.set('rIns', '18288')   # ~0.02in right
    bodyPr.set('tIns', '0')       # zero top
    bodyPr.set('bIns', '0')       # zero bottom

def _lock_row_height(tbl, row_idx, height_emu):
    """Force a fixed row height in the XML so PowerPoint cannot auto-expand it."""
    tr = tbl._tbl.tr_lst[row_idx]
    trPr = tr.find(qn('a:trPr'))
    if trPr is None:
        trPr = etree.SubElement(tr, qn('a:trPr'))
        tr.insert(0, trPr)
    tr.set('h', str(int(height_emu)))

def add_table_parts(prs, policy, df, title, logo_path):
    blank = prs.slide_layouts[6]
    usable_w = prs.slide_width - Inches(2 * policy.side_in)
    TITLE_TOP = Inches(0.65)
    TITLE_H = Inches(0.25)
    TITLE_FONT = Pt(13)
    TABLE_TOP = Inches(0.92)

    # Hard bottom boundary: slide height minus 2 cm
    BOTTOM_MARGIN_CM = 2.0
    max_table_bottom = prs.slide_height - Emu(int(BOTTOM_MARGIN_CM * 360000))
    max_table_h = max_table_bottom - TABLE_TOP

    HEADER_ROW_H = Inches(0.25)
    DATA_ROW_H = Inches(0.22)

    # Compute column widths up front so we can estimate text wrapping
    weights = content_weights(df)
    cols = len(df.columns) if not df.empty else 1
    tw = sum(weights[:cols]) or 1.0
    col_widths_emu = [max(int(usable_w * (weights[c] / tw)), Inches(0.40)) for c in range(cols)] if cols > 0 else [int(usable_w)]
    # Clamp to usable_w
    total_cw = sum(col_widths_emu)
    if total_cw > usable_w:
        scale = usable_w / total_cw
        col_widths_emu = [max(int(w * scale), Inches(0.40)) for w in col_widths_emu]

    # Estimate max text lines per data row to predict PowerPoint row expansion.
    # Use 6pt font ~ 16 chars/inch as baseline.
    CHARS_PER_INCH = 16
    max_lines = 1
    if not df.empty:
        for _, row in df.iterrows():
            for ci in range(cols):
                val = str(row.iloc[ci]) if not pd.isna(row.iloc[ci]) else ''
                col_w_in = col_widths_emu[ci] / 914400.0 if ci < len(col_widths_emu) else 1.0
                chars_fit = max(int(col_w_in * CHARS_PER_INCH), 1)
                lines_needed = max(1, -(-len(val) // chars_fit))
                if lines_needed > max_lines:
                    max_lines = lines_needed

    # Also check header wrapping
    hdr_max_lines = 1
    if not df.empty:
        for ci, col_name in enumerate(df.columns):
            col_w_in = col_widths_emu[ci] / 914400.0 if ci < len(col_widths_emu) else 1.0
            chars_fit = max(int(col_w_in * CHARS_PER_INCH), 1)
            lines_needed = max(1, -(-len(str(col_name)) // chars_fit))
            if lines_needed > hdr_max_lines:
                hdr_max_lines = lines_needed

    effective_hdr_h = max(int(HEADER_ROW_H), int(HEADER_ROW_H * hdr_max_lines))
    effective_row_h = max(int(DATA_ROW_H), int(DATA_ROW_H * max_lines))

    max_data_rows = int((max_table_h - effective_hdr_h) / effective_row_h)
    if max_data_rows < 1:
        max_data_rows = 1
    effective_rows_per_part = min(policy.rows_per_part, max_data_rows)

    parts = chunk_rows(df, effective_rows_per_part)
    total = len(parts); added = 0
    for idx, part in enumerate(parts, start=1):
        s = prs.slides.add_slide(blank)
        display_title = f"{title} (Part {idx}/{total})" if total > 1 else title
        if idx > 1:
            display_title = f"{title} (cont'd — Part {idx}/{total})"
        tb = s.shapes.add_textbox(Inches(policy.side_in), TITLE_TOP, usable_w, TITLE_H)
        tf = tb.text_frame; tf.clear(); tf.word_wrap = False
        tf.text = display_title
        tpf = tf.paragraphs[0]; tpf.font.name = FONT_NAME; tpf.font.size = TITLE_FONT; tpf.font.bold = True; tpf.font.color.rgb = NAVY
        rows = len(part) + 1; part_cols = len(part.columns)
        needed_h = effective_hdr_h + effective_row_h * (rows - 1)
        tbl_h = min(needed_h, int(max_table_h))
        tbl_shape = s.shapes.add_table(rows, part_cols, Inches(policy.side_in), TABLE_TOP, usable_w, tbl_h)
        tbl = tbl_shape.table
        if part_cols > 0 and 'Release' in list(part.columns):
            part = part.copy(); part['Release'] = part['Release'].ffill().fillna('')
        # Use 6pt font for dense tables (8+ cols), otherwise policy font
        hdr_pt = min(policy.header_pt, 6) if part_cols >= 8 else policy.header_pt
        bod_pt = min(policy.body_pt, 6) if part_cols >= 8 else policy.body_pt
        tbl.rows[0].height = effective_hdr_h
        _lock_row_height(tbl, 0, effective_hdr_h)
        for c, col in enumerate(part.columns):
            cell = tbl.cell(0, c); cell.text = str(col)
            cell.text_frame.word_wrap = False
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT_NAME; para.font.size = Pt(hdr_pt); para.font.bold = True; para.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            _set_cell_autofit(cell)
        for r in range(1, rows):
            tbl.rows[r].height = effective_row_h
            _lock_row_height(tbl, r, effective_row_h)
            for c in range(part_cols):
                val = part.iat[r - 1, c]
                if pd.isna(val): val = ''
                cell = tbl.cell(r, c); cell.text = str(val)
                cell.text_frame.word_wrap = False
                para = cell.text_frame.paragraphs[0]; para.font.name = FONT_NAME; para.font.size = Pt(bod_pt)
                _set_cell_autofit(cell)
        # Apply column widths
        for c in range(part_cols):
            tbl.columns[c].width = col_widths_emu[c] if c < len(col_widths_emu) else Inches(0.40)
        if part_cols > 0 and 'Release' in list(part.columns):
            try:
                merge_release_column(tbl, part, list(part.columns).index('Release'), body_font_pt=bod_pt)
            except Exception:
                pass
        add_brand_footer_logo(s, prs, logo_path)
        added += 1
    return added

def add_thank_you(prs, logo_path):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    t = s.shapes.add_textbox(Inches(2.0), Inches(3.2), prs.slide_width - Inches(4.0), Inches(1.0))
    tf = t.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = 'Thank You'; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40); p.font.bold = True; p.font.name = FONT_NAME; p.font.color.rgb = NAVY
    add_brand_footer_logo(s, prs, logo_path)

def _call_automation(func, browser, iteration: str, integration: str):
    params = list(inspect.signature(func).parameters.keys())
    mapping = {
        'browser': browser, 'iteration': iteration, 'iter': iteration,
        'integration': integration, 'integration_str': integration,
        'integration_name': integration, 'project': integration,
    }
    kw = {k: mapping[k] for k in params if k in mapping}
    try:
        return func(**kw)
    except TypeError:
        pass
    candidates = [
        (browser, iteration, integration), (browser, iteration),
        (iteration, integration), (iteration,), (browser,), tuple(),
    ]
    last_err = None
    for args in candidates:
        try:
            return func(*args)
        except TypeError as e:
            last_err = e; continue
    raise last_err if last_err else RuntimeError('Unable to call automation function')

# ---- QMetry Test Report constants & helpers ----
_QTR_HEADLESS = os.getenv('HEADLESS', 'false').lower() == 'true'
_QTR_TIMEOUT_MS = int(os.getenv('TIMEOUT_MS', '45000'))
_QTR_GRID_SETTLE_MS = int(os.getenv('GRID_SETTLE_MS', '2000'))
_QTR_DOWNLOAD_TIMEOUT_MS = int(os.getenv('DOWNLOAD_TIMEOUT_MS', '30000'))
_QTR_PROJECT_ID = os.getenv('QMETRY_PROJECT_ID', '127207')
_QTR_FAST_MODE = os.getenv('FAST_MODE', 'true').lower() == 'true'
_QTR_POST_FILTER_WAIT_MS = int(os.getenv('POST_FILTER_WAIT_MS', '1200' if _QTR_FAST_MODE else '7000'))

_QTR_TBL_BTN_XPATH = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[2]/div/div[1]/div[1]/div/div[4]/div/button/span/span/span"
_QTR_EXP_BTN_XPATH = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[2]/div/div[1]/div[2]/div/div[2]/div/div/div[1]/div/div/div/button/span/span[2]"
_QTR_CSV_ITEM_XPATH = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[2]/div/div[1]/div[2]/div/div[2]/div/div/div[3]/div/div/div/div/div/div/span[3]/span/span/span"

_qtr_t0 = time.time()
def _qtr_tnow(): return int((time.time() - _qtr_t0) * 1000)
def _qtr_info(msg): print(msg, flush=True)

def _qtr_wait_page_idle(page, state='networkidle', timeout=30000):
    try: page.wait_for_load_state(state, timeout=timeout)
    except PWTimeoutError: page.wait_for_timeout(400)

def _qtr_click_xpath(page, xpath, name, wait_after_ms=0):
    loc = page.locator(f'xpath={xpath}').first
    loc.wait_for(state='attached', timeout=_QTR_TIMEOUT_MS)
    try: loc.scroll_into_view_if_needed()
    except Exception: pass
    try: loc.click()
    except Exception: loc.evaluate('el => el.click()')
    if wait_after_ms: page.wait_for_timeout(wait_after_ms)

def _qtr_login_if_needed(page):
    _qtr_wait_page_idle(page)
    if page.locator("input[type='password']").count() > 0:
        if JIRA_USER and JIRA_PASS:
            _qtr_info(f'[{_qtr_tnow()} ms] Login page detected')
            for sel, val in [("input[name='username'], input[type='text']", JIRA_USER), ("input[name='password']", JIRA_PASS)]:
                try: page.fill(sel, val)
                except Exception: pass
            for sel in ["button[type='submit']", "//button[contains(.,'Log in')]"]:
                try: page.click(sel)
                except Exception: pass
            _qtr_wait_page_idle(page)
        else:
            _qtr_info(f'[{_qtr_tnow()} ms] Manual login required')
            try: input()
            except Exception: pass
            _qtr_wait_page_idle(page)

def _qtr_goto_qmetry(page):
    page.wait_for_selector("//a[contains(text(), 'More')]")
    page.locator("//a[contains(text(), 'More')]").first.click()
    page.wait_for_timeout(500)
    page.locator("//a[contains(text(), 'QMetry')]").first.hover()
    page.wait_for_timeout(400)
    page.locator("//a[contains(text(), 'Test Management')]").first.click()
    _qtr_wait_page_idle(page); page.wait_for_timeout(600)

def _qtr_select_software_project(page):
    page.locator("//span[contains(text(), 'Software Project')]").first.click()
    page.wait_for_timeout(300)
    page.keyboard.type('MDA-Integration'); page.keyboard.press('Enter')
    page.wait_for_timeout(500)

def _qtr_search_integration_folder(page, iteration_str):
    page.locator('div.glyph-search').first.click()
    page.wait_for_timeout(200)
    page.keyboard.type(f'INTEGRATION_PROGRESSION_{iteration_str}')
    page.keyboard.press('Enter')
    _qtr_wait_page_idle(page)
    page.wait_for_timeout(_QTR_POST_FILTER_WAIT_MS)
    _qtr_info(f'[{_qtr_tnow()} ms] Left filter applied; waited {_QTR_POST_FILTER_WAIT_MS} ms')

def _qtr_is_on_test_cycle(page):
    try: return '#/Manage/TestCycle' in (page.url or '')
    except Exception: return False

def _qtr_go_to_test_cycle(page):
    page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/Manage/TestCycle?projectId={_QTR_PROJECT_ID}')
    _qtr_wait_page_idle(page); page.wait_for_timeout(400)
    if not _qtr_is_on_test_cycle(page):
        for xp in [
            "//*[@id='root_4.14.1.1']//*[@role='tab' and contains(normalize-space(),'Test Cycle')]",
            "//*[@id='root_4.14.1.1']//*[contains(@class,'react-tabs__tab') and contains(normalize-space(),'Test Cycle')]",
            "//*[@id='root_4.14.1.1']//*[self::span or self::div or self::button][contains(normalize-space(),'Test Cycle')]",
        ]:
            try:
                _qtr_click_xpath(page, xp, 'Test Cycle', 300)
                if _qtr_is_on_test_cycle(page): break
            except Exception: continue
    if not _qtr_is_on_test_cycle(page):
        raise RuntimeError("Unable to open 'Test Cycle'.")
    _qtr_info(f'[{_qtr_tnow()} ms] Test Cycle active')

_QTR_SPINNER_SELECTORS = [
    'div.spinner', 'div.loader', 'div.page-loader', 'div.custom-loader',
    "svg[aria-label='Loading']", "[role='progressbar']",
    "div.ReactVirtualized__Grid__innerScrollContainer div[style*='translateY']"
]
_QTR_HEADER_TOKENS = ['Summary', 'Key', 'Execution Result']

def _qtr_frame_candidates(page):
    frames = [page]
    for fr in page.frames:
        if fr != page.main_frame: frames.append(fr)
    return frames

def _qtr_wait_for_test_cycle_grid_stable(page):
    start = _qtr_tnow(); deadline = time.time() + (_QTR_TIMEOUT_MS / 1000)
    header_ok = False
    while time.time() < deadline and not header_ok:
        for ctx in _qtr_frame_candidates(page):
            for token in _QTR_HEADER_TOKENS:
                try:
                    if ctx.locator(f'text={token}').count() > 0: header_ok = True; break
                except Exception: pass
            if header_ok: break
        if not header_ok: page.wait_for_timeout(150)
    while time.time() < deadline:
        visible = False
        for sel in _QTR_SPINNER_SELECTORS:
            try:
                loc = page.locator(sel)
                if loc.count() > 0 and loc.first.is_visible(): visible = True; break
            except Exception: pass
        if not visible: break
        page.wait_for_timeout(120)
    _qtr_info(f'[{_qtr_tnow()} ms] Test Cycle grid headers ready in {_qtr_tnow() - start} ms')

_qtr_key_pat = re.compile(r'[A-Z]+-[A-Z]+-\d+')

def _qtr_find_key_in_context(ctx, search_text):
    try: ctx.locator('text=Summary').first.wait_for(timeout=1500)
    except Exception: pass
    row = ctx.locator('table#listingTable tr', has_text=search_text).first
    if row.count() == 0: row = ctx.locator('tr', has_text=search_text).first
    if row.count() == 0: row = ctx.locator('div.qtmTable__row', has_text=search_text).first
    if row.count() == 0:
        cell = ctx.locator(f"xpath=//*[contains(normalize-space(), {repr(search_text)})]").first
        if cell.count() > 0:
            row = cell.locator('xpath=ancestor::tr[1]').first
            if row.count() == 0:
                row = cell.locator("xpath=ancestor::*[contains(@class,'qtmTable__row')][1]").first
    if row.count() == 0: return None, None
    try:
        key_anchor = row.locator('a').filter(has_text=_qtr_key_pat)
        if key_anchor.count() == 0: key_anchor = row.locator('a').first
        key_anchor.wait_for(timeout=2000)
        key_text = key_anchor.inner_text().strip()
        if key_text: return key_text, row
    except Exception: pass
    return None, None

def _qtr_fetch_key_for_summary(page, iteration_str):
    start = _qtr_tnow()
    term_primary = f'TMO_INTEGRATION_PROGRESSION_{iteration_str}'
    term_secondary = f'INTEGRATION_PROGRESSION_{iteration_str}'
    _qtr_wait_for_test_cycle_grid_stable(page)
    def _try(search_text):
        key_text, row = _qtr_find_key_in_context(page, search_text)
        if not key_text:
            for fr in page.frames:
                if fr == page.main_frame: continue
                try:
                    key_text, row = _qtr_find_key_in_context(fr, search_text)
                    if key_text: break
                except Exception: continue
        if key_text and row:
            _qtr_info(f"[{_qtr_tnow()} ms] Match term: '{search_text}'  Key: {key_text}")
        return key_text
    key = _try(term_primary) or _try(term_secondary)
    if not key:
        raise RuntimeError('Summary row not found for either candidate.')
    _qtr_info(f'[{_qtr_tnow()} ms] Key fetched in {_qtr_tnow() - start} ms')
    return key

def _qtr_go_to_test_report(page):
    page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/Manage/TestReport?projectId={_QTR_PROJECT_ID}')
    _qtr_wait_page_idle(page); page.wait_for_timeout(500)
    try:
        page.locator("xpath=//*[@id='root_4.14.1.1']//span[normalize-space()='Test Execution Reports']").first.wait_for(timeout=_QTR_TIMEOUT_MS)
    except Exception:
        try: _qtr_click_xpath(page, "//*[@id='root_4.14.1.1']/div/div/div/div[1]/div[1]/div[4]/div[4]", 'Test Report tab', 500)
        except Exception: pass
    for xp in [
        "//*[@id='root_4.14.1.1']//span[normalize-space()='Test Execution Reports']",
        "//*[@id='root_4.14.1.1']//*[normalize-space()='Test Case Execution Summary']",
        "//*[@id='root_4.14.1.1']//*[normalize-space()='By Execution Assignee']",
    ]:
        try: _qtr_click_xpath(page, xp, 'Left tree select', 300)
        except Exception: continue
    _qtr_info(f'[{_qtr_tnow()} ms] Test Report > By Execution Assignee ready')

def _qtr_filter_by_key_and_generate(page, key_value):
    start = _qtr_tnow()
    _qtr_wait_page_idle(page); page.wait_for_timeout(800)
    input_xpath = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[3]/div[4]/div/div[2]/div/div[2]/div[2]/div/input"
    field = page.locator(f'xpath={input_xpath}').first
    field.wait_for(timeout=_QTR_TIMEOUT_MS)
    try: field.fill('')
    except Exception:
        field.click(); page.keyboard.press('Control+A'); page.keyboard.press('Delete')
    field.type(key_value, delay=15)
    page.wait_for_timeout(600); page.keyboard.press('Enter'); page.wait_for_timeout(900)
    gen_btn_xpath = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[3]/div[8]/div/div/div/button/span/span"
    page.locator(f'xpath={gen_btn_xpath}').first.click(); page.wait_for_timeout(1200)
    _qtr_info(f'[{_qtr_tnow()} ms] Filter & Generate took {_qtr_tnow() - start} ms')

def _qtr_export_csv_tabular(page, context):
    step_start = _qtr_tnow()
    try: _qtr_click_xpath(page, _QTR_TBL_BTN_XPATH, 'Tabular icon', wait_after_ms=0)
    except Exception:
        for xp in [
            "//button[.//span[normalize-space()='Tabular View'] or normalize-space()='Tabular View']",
            "//button//*[contains(@class,'icon') and (contains(@class,'table') or contains(@class,'grid'))]/ancestor::button[1]",
        ]:
            try: _qtr_click_xpath(page, xp, 'Tabular (fallback)', wait_after_ms=0); break
            except Exception: continue
    page.wait_for_timeout(3000)
    try: _qtr_click_xpath(page, _QTR_EXP_BTN_XPATH, 'Export button', wait_after_ms=0)
    except Exception:
        for xp in [
            "//button[.//span[normalize-space()='Export']]",
            "//button[normalize-space()='Export']",
            "//*[self::span or self::a or self::div][normalize-space()='Export']/ancestor::button[1]",
        ]:
            try: _qtr_click_xpath(page, xp, 'Export button (fallback)', wait_after_ms=0); break
            except Exception: continue
    page.wait_for_timeout(1000)
    with page.expect_download(timeout=_QTR_DOWNLOAD_TIMEOUT_MS) as dl_info:
        clicked = False
        try: _qtr_click_xpath(page, _QTR_CSV_ITEM_XPATH, 'CSV with Tabular Summary', wait_after_ms=0); clicked = True
        except Exception: pass
        if not clicked:
            for xp in [
                "//*[@role='menu']//*[normalize-space()='CSV with Tabular Summary']",
                "//*[normalize-space()='CSV with Tabular Summary']",
                "//span[normalize-space()='CSV with Tabular Summary']",
            ]:
                try: _qtr_click_xpath(page, xp, 'CSV (text fallback)', wait_after_ms=0); clicked = True; break
                except Exception: continue
        if not clicked:
            raise RuntimeError('CSV with Tabular Summary option not found in dropdown')
    download = dl_info.value
    suggested = download.suggested_filename or 'export.csv'
    target_path = DOWNLOADS_DIR / suggested
    download.save_as(target_path)
    deadline = time.time() + (_QTR_DOWNLOAD_TIMEOUT_MS / 1000)
    while time.time() < deadline:
        if target_path.exists() and target_path.stat().st_size > 0:
            _qtr_info(f'[{_qtr_tnow()} ms] Export completed in {_qtr_tnow() - step_start} ms\n \u2022 Downloaded file: {target_path}')
            return target_path
        time.sleep(0.2)
    raise TimeoutError(f'Download not found or zero-sized after {_QTR_DOWNLOAD_TIMEOUT_MS} ms: {target_path}')

_QTR_STATUS_CANDIDATES = [
    'BLOCKED', 'FAIL', 'WORK IN PROGRESS', 'NOT EXECUTED', 'PASS',
    'IN PROGRESS', 'NOT APPLICABLE', 'CONDITIONAL PASS', 'DEFERRED'
]
_qtr_numeric_cell_pat = re.compile(r'^\s*(\d+)\s*(?:\(.*?%\))?\s*$')

def _qtr_as_int(v):
    if pd.isna(v): return 0
    if isinstance(v, (int, float)):
        try: return int(v)
        except Exception: return 0
    s = str(v).strip()
    if s in ('', '-', '\u2013', '\u2014', None): return 0
    m = _qtr_numeric_cell_pat.match(s)
    if m:
        try: return int(m.group(1))
        except Exception: return 0
    try: return int(float(s))
    except Exception: return 0

def _qtr_post_process_csv_to_excel(csv_path, iteration_str):
    try: df = pd.read_csv(csv_path)
    except Exception: df = pd.read_csv(csv_path, sep=';')
    df.columns = [str(c).strip() for c in df.columns]
    exec_candidates = ['Execution Assignee', 'Assignee', 'Execution_Assignee', 'execution assignee', 'execution_assignee']
    col_map = {c.lower(): c for c in df.columns}
    exec_col = next((col_map[c.lower()] for c in exec_candidates if c.lower() in col_map), None)
    if exec_col is None:
        _qtr_info(f"[WARN] 'Execution Assignee' column not found. Columns: {list(df.columns)}. Skipping filter.")
        df_f = df.copy()
    else:
        ser = df[exec_col].astype(str).str.strip()
        df_f = df[ser.isin(ALLOWED_REPORTERS)].copy()
        _qtr_info(f'[{_qtr_tnow()} ms] Filtered by allowed assignees: kept {len(df_f)} of {len(df)} rows')
    present_status = [c for c in _QTR_STATUS_CANDIDATES if c in df_f.columns]
    df_f.insert(0, 'Release', f'INTEGRATION_PROGRESSION_{iteration_str}')
    ordered_cols = ['Release'] + ([exec_col] if exec_col else []) + present_status
    try: df_f = df_f[ordered_cols]
    except Exception: pass
    for c in present_status:
        df_f[c] = df_f[c].map(_qtr_as_int)
    if present_status:
        df_f['TOTAL'] = df_f[present_status].sum(axis=1)
    if present_status:
        totals_dict = {c: int(df_f[c].sum()) for c in present_status}
        totals_dict['TOTAL'] = int(df_f['TOTAL'].sum()) if 'TOTAL' in df_f.columns else sum(totals_dict.values())
        totals_dict['Release'] = f'INTEGRATION_PROGRESSION_{iteration_str}'
        if exec_col: totals_dict[exec_col] = 'TOTAL'
        total_row = {col: totals_dict.get(col, '') for col in df_f.columns}
        df_out = pd.concat([df_f, pd.DataFrame([total_row])], ignore_index=True)
    else:
        df_out = df_f
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    out_path = csv_path.with_name(f'Test_Summary_report_{ts}.xlsx')
    with pd.ExcelWriter(out_path, engine='openpyxl') as writer:
        df_out.to_excel(writer, index=False)
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Font
        from openpyxl.utils import get_column_letter
        wb = load_workbook(out_path); ws = wb.active
        last_row = ws.max_row; bold_font = Font(bold=True)
        for c in range(1, ws.max_column + 1):
            ws.cell(row=last_row, column=c).font = bold_font
        ws.freeze_panes = 'A2'
        for col_idx in range(1, ws.max_column + 1):
            col_letter = get_column_letter(col_idx)
            max_len = 0
            for row_idx in range(1, ws.max_row + 1):
                val = ws.cell(row=row_idx, column=col_idx).value
                max_len = max(max_len, len(str(val)) if val is not None else 0)
            ws.column_dimensions[col_letter].width = min(max(8, int(max_len * 1.2) + 2), 60)
        wb.save(out_path)
        _qtr_info(f'[{_qtr_tnow()} ms] Styled TOTAL + freeze header + auto-fit widths in Excel')
    except Exception as e:
        _qtr_info(f'[WARN] Could not apply Excel formatting: {e}')
    return out_path

def _qtr_post_process_label_csv(csv_path, iteration_str):
    """Process the By Label CSV into a Features in Scope Excel matching the expected layout."""
    try:
        df = pd.read_csv(csv_path)
    except Exception:
        df = pd.read_csv(csv_path, sep=';')
    df.columns = [str(c).strip() for c in df.columns]
    _qtr_info(f'[{_qtr_tnow()} ms] Label CSV: {len(df)} rows, cols={list(df.columns)}')

    # Find the label column
    label_col = None
    for c in df.columns:
        if 'label' in c.lower():
            label_col = c; break
    if label_col is None:
        label_col = df.columns[0]

    # Map status columns to integers
    status_cols = ['BLOCKED', 'FAIL', 'WORK IN PROGRESS', 'NOT EXECUTED', 'PASS',
                   'IN PROGRESS', 'NOT APPLICABLE', 'CONDITIONAL PASS', 'DEFERRED']
    present = [c for c in status_cols if c in df.columns]
    for c in present:
        df[c] = df[c].apply(_qtr_as_int)

    # Helper to get column value or 0
    def _col_val(row, *names):
        for n in names:
            if n in row.index:
                return int(row[n]) if not pd.isna(row[n]) else 0
        return 0

    # Build Features table matching expected layout
    rows_out = []
    for _, row in df.iterrows():
        passed = _col_val(row, 'PASS')
        failed = _col_val(row, 'FAIL')
        blocked = _col_val(row, 'BLOCKED')
        in_prog = _col_val(row, 'IN PROGRESS', 'WORK IN PROGRESS')
        not_exec = _col_val(row, 'NOT EXECUTED')
        no_run = _col_val(row, 'NOT APPLICABLE', 'DEFERRED', 'CONDITIONAL PASS')
        total = passed + failed + blocked + in_prog + not_exec + no_run
        pass_pct = f"{int(passed / total * 100)}%" if total > 0 else '0%'
        fail_pct = f"{int(failed / total * 100)}%" if total > 0 else '0%'
        exec_pct = f"{int((total - not_exec) / total * 100)}%" if total > 0 else '0%'
        test_status = 'Test Complete' if not_exec == 0 and in_prog == 0 else 'In Progress'
        rows_out.append({
            'Feature Id': row[label_col],
            'Test Status': test_status,
            'Passed': passed,
            'Pass %': pass_pct,
            'Failed': failed,
            'Blocked': blocked,
            'InProgress': in_prog,
            'Fail %': fail_pct,
            'No Run': no_run,
            'Total': total,
            'Execution %': exec_pct,
        })

    features = pd.DataFrame(rows_out)

    # Add Total row at top (like the screenshot)
    t_passed = int(features['Passed'].sum())
    t_failed = int(features['Failed'].sum())
    t_blocked = int(features['Blocked'].sum())
    t_inprog = int(features['InProgress'].sum())
    t_norun = int(features['No Run'].sum())
    t_total = int(features['Total'].sum())
    total_row = {
        'Feature Id': 'Total',
        'Test Status': '',
        'Passed': t_passed,
        'Pass %': f"{int(t_passed / t_total * 100)}%" if t_total > 0 else '0%',
        'Failed': t_failed,
        'Blocked': t_blocked,
        'InProgress': t_inprog,
        'Fail %': f"{int(t_failed / t_total * 100)}%" if t_total > 0 else '0%',
        'No Run': t_norun,
        'Total': t_total,
        'Execution %': f"{int((t_total - int(features[features['Feature Id'] != 'Total']['No Run'].sum()) - sum(_col_val(r, 'NOT EXECUTED') for _, r in df.iterrows())) / t_total * 100)}%" if t_total > 0 else '0%',
    }
    # Simpler execution % for total
    t_not_exec = sum(_col_val(r, 'NOT EXECUTED') for _, r in df.iterrows())
    total_row['Execution %'] = f"{int((t_total - t_not_exec) / t_total * 100)}%" if t_total > 0 else '0%'
    features = pd.concat([pd.DataFrame([total_row]), features], ignore_index=True)

    # Save
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    out_path = ARTIFACTS / f'Features_in_Scope_{iteration_str}_{ts}.xlsx'
    features.to_excel(out_path, index=False)
    _qtr_info(f'[{_qtr_tnow()} ms] Features in Scope saved: {out_path} ({len(features)} rows)')
    return out_path


def run_qmetry_test_report(iteration, keep_browser=False):
    if keep_browser:
        pw = sync_playwright().start()
        browser = pw.chromium.launch(headless=_QTR_HEADLESS, channel='msedge')
        context = browser.new_context(accept_downloads=True)
        page = context.new_page()
    else:
        pw = None
    try:
        if not keep_browser:
            _pw_cm = sync_playwright()
            _p = _pw_cm.__enter__()
            browser = _p.chromium.launch(headless=_QTR_HEADLESS, channel='msedge')
            context = browser.new_context(accept_downloads=True)
            page = context.new_page()

        _qtr_info(f'[{_qtr_tnow()} ms] === QMetry Test Report (embedded) (iteration={iteration}) ===')
        page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/?projectId={_QTR_PROJECT_ID}')
        _qtr_wait_page_idle(page)
        _qtr_login_if_needed(page)
        _qtr_goto_qmetry(page)
        _qtr_select_software_project(page)
        _qtr_search_integration_folder(page, iteration)
        _qtr_go_to_test_cycle(page)
        key_value = _qtr_fetch_key_for_summary(page, iteration)
        _qtr_go_to_test_report(page)
        _qtr_filter_by_key_and_generate(page, key_value)

        # ── Step A: By Execution Assignee export ──
        excel_path = None
        try:
            csv_path = _qtr_export_csv_tabular(page, context)
            excel_path = _qtr_post_process_csv_to_excel(csv_path, iteration)
            _qtr_info(f'[{_qtr_tnow()} ms] === By Assignee completed === Excel: {excel_path}')
        except Exception as assignee_err:
            _qtr_info(f'[WARN] By Assignee export failed: {assignee_err}')

        # ── Step B: By Label export ──
        label_excel_path = None
        try:
            _qtr_info(f'[{_qtr_tnow()} ms] === Step: Switching to By Label ===')
            for xp in [
                "//*[@id='root_4.14.1.1']//span[normalize-space()='Test Execution Reports']",
                "//*[@id='root_4.14.1.1']//*[normalize-space()='Test Case Execution Summary']",
            ]:
                try: _qtr_click_xpath(page, xp, 'Expand tree', 300)
                except Exception: continue
            page.wait_for_timeout(500)
            _BY_LABEL_XPATH = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[1]/nav/div/div[2]/div/ul/li[4]/ul/li[2]/ul/li[6]/div/div/div/span"
            loc = page.locator(f'xpath={_BY_LABEL_XPATH}').first
            loc.wait_for(state='attached', timeout=10000)
            try: loc.scroll_into_view_if_needed()
            except Exception: pass
            loc.click()
            _qtr_info(f'[{_qtr_tnow()} ms] === Clicked By Label ===')
            page.wait_for_timeout(1500)
            _qtr_filter_by_key_and_generate(page, key_value)
            label_csv_path = _qtr_export_csv_tabular(page, context)
            label_excel_path = _qtr_post_process_label_csv(label_csv_path, iteration)
            _qtr_info(f'[{_qtr_tnow()} ms] === By Label completed === Excel: {label_excel_path}')
        except Exception as lbl_err:
            _qtr_info(f'[WARN] By Label failed: {lbl_err}')
            traceback.print_exc()

        if keep_browser:
            return excel_path, label_excel_path, page, context, browser, pw
        return excel_path, label_excel_path
    except Exception as e:
        _qtr_info(f'[ERROR] run_qmetry_test_report FAILED: {e}')
        traceback.print_exc()
        if keep_browser:
            return None, None, page, context, browser, pw
        raise
    finally:
        if not keep_browser:
            try: context.close()
            except Exception: pass
            try: browser.close()
            except Exception: pass
            try: _pw_cm.__exit__(None, None, None)
            except Exception: pass


# ═══════════════════════════════════════════════════════════════════════════════
# V3.0: QMetry Test Cycle Pie Chart Reports
# ═══════════════════════════════════════════════════════════════════════════════

def _qtr_fetch_both_keys(page, iteration_str):
    """Fetch both Progression and Regression keys from the Test Cycle grid."""
    _qtr_wait_for_test_cycle_grid_stable(page)
    def _try_find(search_text):
        key_text, row = _qtr_find_key_in_context(page, search_text)
        if not key_text:
            for fr in page.frames:
                if fr == page.main_frame: continue
                try:
                    key_text, row = _qtr_find_key_in_context(fr, search_text)
                    if key_text: break
                except Exception: continue
        if key_text:
            _qtr_info(f"[{_qtr_tnow()} ms] Found key for '{search_text}': {key_text}")
        return key_text
    prog_key = _try_find(f'TMO_INTEGRATION_PROGRESSION_{iteration_str}')
    reg_key = _try_find(f'TMO_INTEGRATION_REGRESSION_{iteration_str}')
    if not prog_key:
        prog_key = _try_find(f'INTEGRATION_PROGRESSION_{iteration_str}')
    if not reg_key:
        reg_key = _try_find(f'INTEGRATION_REGRESSION_{iteration_str}')
    _qtr_info(f'[{_qtr_tnow()} ms] Keys: progression={prog_key}, regression={reg_key}')
    return prog_key, reg_key


def _cycle_scrape_all_tables(page):
    """Scrape all visible tables from the page and its frames."""
    all_dfs = []
    for ctx in _qtr_frame_candidates(page):
        try:
            tables_html = ctx.evaluate("""() => {
                const tables = document.querySelectorAll('table');
                return Array.from(tables).map(t => t.outerHTML);
            }""")
            for html in tables_html:
                try:
                    dfs = pd.read_html(io.StringIO(html))
                    if dfs and len(dfs[0]) > 0:
                        all_dfs.append(dfs[0])
                except Exception:
                    pass
        except Exception:
            pass
    return all_dfs


def run_qmetry_cycle_reports(iteration, page=None):
    """Navigate to QMetry Reports > By Test Cycle, capture pie chart screenshots."""
    results = {'progression_pie_img': None, 'regression_pie_img': None, 'features_df': None}
    own_browser = page is None
    p_ctx = None; browser = None; context = None
    try:
        if own_browser:
            p_ctx = sync_playwright().start()
            browser = p_ctx.chromium.launch(headless=_QTR_HEADLESS, channel='msedge')
            context = browser.new_context(accept_downloads=True, viewport={'width': 1920, 'height': 1080})
            page = context.new_page()
            page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/?projectId={_QTR_PROJECT_ID}',
                       timeout=120000, wait_until='domcontentloaded')
            _qtr_wait_page_idle(page)
            _qtr_login_if_needed(page)
            _qtr_goto_qmetry(page)
            _qtr_select_software_project(page)
            _qtr_search_integration_folder(page, iteration)
            _qtr_go_to_test_cycle(page)
            _qtr_wait_for_test_cycle_grid_stable(page)
        else:
            _qtr_info(f'[{_qtr_tnow()} ms] === QMetry Cycle Reports (reusing session) ===')
            _qtr_go_to_test_cycle(page)
            _qtr_wait_page_idle(page)
            page.wait_for_timeout(3000)
            # Re-select project and search for the folder
            try:
                _qtr_select_software_project(page)
            except Exception:
                pass
            try:
                _qtr_search_integration_folder(page, iteration)
            except Exception:
                pass
            page.wait_for_timeout(3000)
            _qtr_wait_for_test_cycle_grid_stable(page)
            page.wait_for_timeout(2000)

        # Step 1: Fetch both keys (with retry)
        prog_key, reg_key = _qtr_fetch_both_keys(page, iteration)
        if not prog_key and not reg_key:
            # Retry: wait longer and try again
            _qtr_info(f'[{_qtr_tnow()} ms] Keys not found, retrying after wait...')
            page.wait_for_timeout(5000)
            _qtr_wait_for_test_cycle_grid_stable(page)
            prog_key, reg_key = _qtr_fetch_both_keys(page, iteration)
        if not prog_key and not reg_key:
            _qtr_info(f'[WARN] Keys still not found — skipping cycle reports')
            return results
        keys_csv = ','.join(k for k in [prog_key, reg_key] if k)
        _qtr_info(f'[{_qtr_tnow()} ms] Keys for filter: {keys_csv}')

        # Step 2: Navigate to Reports > By Test Cycle
        page.goto(f'{JIRA_BASE_URL}/secure/QTMAction.jspa#/Manage/TestReport?projectId={_QTR_PROJECT_ID}')
        _qtr_wait_page_idle(page); page.wait_for_timeout(1000)
        for xp in [
            "//*[@id='root_4.14.1.1']//span[normalize-space()='Test Execution Reports']",
            "//*[@id='root_4.14.1.1']//*[normalize-space()='Test Case Execution Summary']",
            "//*[@id='root_4.14.1.1']//*[normalize-space()='By Test Cycle']",
        ]:
            try: _qtr_click_xpath(page, xp, 'Left tree', 500)
            except Exception: continue
        page.wait_for_timeout(1500)

        # Step 3: Enter keys and Generate
        _qtr_wait_page_idle(page); page.wait_for_timeout(800)
        input_xp = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[3]/div[4]/div/div[2]/div/div[2]/div[2]/div/input"
        field = page.locator(f'xpath={input_xp}').first
        try: field.wait_for(timeout=_QTR_TIMEOUT_MS)
        except Exception:
            field = page.locator("input[type='text']").last
            field.wait_for(timeout=10000)
        try: field.fill('')
        except Exception:
            field.click(); page.keyboard.press('Control+A'); page.keyboard.press('Delete')
        field.type(keys_csv, delay=15)
        page.wait_for_timeout(600); page.keyboard.press('Enter'); page.wait_for_timeout(900)
        gen_xp = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[3]/div[8]/div/div/div/button/span/span"
        try: page.locator(f'xpath={gen_xp}').first.click()
        except Exception:
            for sel in ["button:has-text('Generate')", "//button[contains(.,'Generate')]"]:
                try: page.locator(sel).first.click(); break
                except Exception: continue
        page.wait_for_timeout(3000)

        # Step 4: Click Pie chart icon (3rd icon in toolbar)
        pie_xp = "//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[2]/div/div[1]/div[1]/div/div[3]/div/button"
        try: page.locator(f'xpath={pie_xp}').first.click()
        except Exception:
            btns = page.locator("//*[@id='root_4.14.1.1']/div/div/div/div[2]/div[2]/div/div/div[2]/div/div[1]/div[1]/div/div/div/button")
            for i in range(btns.count()):
                if i == 2:
                    try: btns.nth(i).click()
                    except Exception: pass
                    break
        page.wait_for_timeout(3000)

        # Step 5: Screenshot pie for each key via dropdown
        def _screenshot_pie(key_val, label):
            if not key_val: return None
            _qtr_info(f'[{_qtr_tnow()} ms] Selecting dropdown for {label}: {key_val}')
            dd_clicked = False

            # The dropdown is a custom React component (not native <select>).
            # Strategy: find any clickable element containing "QMMDA-TR-" text,
            # click it to open the dropdown, then click the option with our key.

            # First, find the dropdown trigger by looking for elements with QMMDA text
            try:
                # Look for the dropdown trigger — it contains text like "QMMDA-TR-XX (TMO_..."
                triggers = page.locator("//*[@id='root_4.14.1.1']//*[contains(text(),'QMMDA-TR-')]")
                _qtr_info(f'[{_qtr_tnow()} ms] Found {triggers.count()} elements with QMMDA text')
                for ti in range(triggers.count()):
                    el = triggers.nth(ti)
                    try:
                        txt = el.inner_text().strip()[:80]
                        tag = el.evaluate('el => el.tagName')
                        _qtr_info(f'[{_qtr_tnow()} ms]   Element #{ti}: <{tag}> "{txt}"')
                    except Exception:
                        pass
            except Exception as de:
                _qtr_info(f'[{_qtr_tnow()} ms] Debug scan error: {de}')

            # Click the dropdown trigger to open it
            try:
                # The dropdown trigger is typically the visible element showing current selection
                dd_trigger = page.locator("//*[@id='root_4.14.1.1']//*[contains(text(),'QMMDA-TR-')]").last
                if dd_trigger.count() > 0 and dd_trigger.is_visible():
                    dd_trigger.click()
                    page.wait_for_timeout(1000)
                    _qtr_info(f'[{_qtr_tnow()} ms] Clicked dropdown trigger')

                    # Now look for the option containing our key in any newly appeared elements
                    # The dropdown options might appear as a list/menu
                    option = page.locator(f"//*[contains(text(),'{key_val}')]").last
                    if option.count() > 0:
                        option.click()
                        dd_clicked = True
                        page.wait_for_timeout(3000)
                        _qtr_info(f'[{_qtr_tnow()} ms] Clicked option: {key_val}')
                    else:
                        _qtr_info(f'[{_qtr_tnow()} ms] Option {key_val} not found in dropdown')
                        page.keyboard.press('Escape')
            except Exception as e:
                _qtr_info(f'[{_qtr_tnow()} ms] Dropdown click error: {e}')

            if not dd_clicked:
                _qtr_info(f'[WARN] Could not switch dropdown to {key_val}')

            # Screenshot just the pie chart area (not full page)
            out_path = ARTIFACTS / f'cycle_pie_{label}_{iteration}.png'
            page.wait_for_timeout(1500)

            # Clip tighter: just the pie chart + labels + legend
            page.screenshot(path=str(out_path), clip={
                'x': 350, 'y': 80, 'width': 1100, 'height': 650
            })
            _qtr_info(f'[{_qtr_tnow()} ms] Pie chart clipped screenshot saved: {out_path}')

            _qtr_info(f'[{_qtr_tnow()} ms] Pie screenshot saved: {out_path}')
            return out_path

        if prog_key:
            results['progression_pie_img'] = _screenshot_pie(prog_key, 'progression')
        if reg_key:
            results['regression_pie_img'] = _screenshot_pie(reg_key, 'regression')
        if not results['progression_pie_img'] and not results['regression_pie_img']:
            fb = ARTIFACTS / f'cycle_pie_fallback_{iteration}.png'
            page.screenshot(path=str(fb), full_page=False)
            results['progression_pie_img'] = fb

        _qtr_info(f'[{_qtr_tnow()} ms] === QMetry Cycle Reports completed ===')
    except Exception as e:
        _qtr_info(f'[ERROR] QMetry Cycle Reports failed: {e}')
        traceback.print_exc()
        try: page.screenshot(path=str(ARTIFACTS / f'cycle_error_{iteration}.png'))
        except Exception: pass
    finally:
        if own_browser:
            try: context.close()
            except Exception: pass
            try: browser.close()
            except Exception: pass
            try: p_ctx.stop()
            except Exception: pass
    return results


def add_screenshot_slide(prs, img_path, title, logo_path):
    """Add a slide with a screenshot image, scaled to fit."""
    if not img_path or not Path(img_path).exists(): return
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    SIDE = Inches(0.60); usable_w = prs.slide_width - Inches(1.20)
    tb = s.shapes.add_textbox(SIDE, Inches(0.40), usable_w, Inches(0.30))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = False; tf.text = title
    p = tf.paragraphs[0]; p.font.name = FONT_NAME; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY
    IMG_TOP = Inches(0.80)
    max_w = prs.slide_width - Inches(1.20)
    max_h = prs.slide_height - IMG_TOP - Inches(0.80)
    from PIL import Image as PILImage
    with PILImage.open(str(img_path)) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    if max_w / max_h > aspect: h = max_h; w = int(h * aspect)
    else: w = max_w; h = int(w / aspect)
    left = int((prs.slide_width - w) / 2)
    s.shapes.add_picture(str(img_path), left, IMG_TOP, w, h)
    add_brand_footer_logo(s, prs, logo_path)


def _color_code_features_table(tbl, df):
    """Apply conditional color coding to a Features in Scope table.
    Green for high Pass%, yellow/amber for medium, red-ish for Fail%."""
    if df is None or df.empty:
        return

    # Color definitions
    TEAL_BG = RGBColor(0, 128, 128)       # Total row / header
    GREEN_100 = RGBColor(198, 224, 180)    # 100% pass
    GREEN_HIGH = RGBColor(216, 235, 200)   # 90-99%
    GREEN_MED = RGBColor(235, 245, 220)    # 70-89%
    YELLOW = RGBColor(255, 242, 204)       # 50-69%
    AMBER = RGBColor(255, 230, 180)        # 30-49%
    RED_LIGHT = RGBColor(255, 210, 210)    # Fail% > 0
    GRAY_LIGHT = RGBColor(242, 242, 242)   # alternating rows

    def _pct_to_int(val):
        s = str(val).strip().replace('%', '')
        try: return int(float(s))
        except Exception: return -1

    def _set_cell_bg(cell, color):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    cols = list(df.columns)
    pass_pct_idx = cols.index('Pass %') if 'Pass %' in cols else -1
    fail_pct_idx = cols.index('Fail %') if 'Fail %' in cols else -1
    exec_pct_idx = cols.index('Execution %') if 'Execution %' in cols else -1
    status_idx = cols.index('Test Status') if 'Test Status' in cols else -1

    for ri in range(1, len(tbl.rows)):  # skip header
        data_ri = ri - 1
        if data_ri >= len(df):
            break
        row_data = df.iloc[data_ri]
        is_total = str(row_data.iloc[0]).strip() == 'Total'

        # Total row gets teal background
        if is_total:
            for ci in range(len(tbl.columns)):
                _set_cell_bg(tbl.cell(ri, ci), TEAL_BG)
                para = tbl.cell(ri, ci).text_frame.paragraphs[0]
                para.font.color.rgb = WHITE
                para.font.bold = True
            continue

        # Alternating row background
        if ri % 2 == 0:
            for ci in range(len(tbl.columns)):
                _set_cell_bg(tbl.cell(ri, ci), GRAY_LIGHT)

        # Color code Pass %
        if pass_pct_idx >= 0 and pass_pct_idx < len(tbl.columns):
            pct = _pct_to_int(row_data.iloc[pass_pct_idx] if pass_pct_idx < len(row_data) else '')
            if pct >= 100:
                _set_cell_bg(tbl.cell(ri, pass_pct_idx), GREEN_100)
            elif pct >= 90:
                _set_cell_bg(tbl.cell(ri, pass_pct_idx), GREEN_HIGH)
            elif pct >= 70:
                _set_cell_bg(tbl.cell(ri, pass_pct_idx), GREEN_MED)
            elif pct >= 50:
                _set_cell_bg(tbl.cell(ri, pass_pct_idx), YELLOW)
            elif pct >= 0:
                _set_cell_bg(tbl.cell(ri, pass_pct_idx), AMBER)

        # Color code Fail %
        if fail_pct_idx >= 0 and fail_pct_idx < len(tbl.columns):
            pct = _pct_to_int(row_data.iloc[fail_pct_idx] if fail_pct_idx < len(row_data) else '')
            if pct > 0:
                _set_cell_bg(tbl.cell(ri, fail_pct_idx), RED_LIGHT)

        # Color code Execution %
        if exec_pct_idx >= 0 and exec_pct_idx < len(tbl.columns):
            pct = _pct_to_int(row_data.iloc[exec_pct_idx] if exec_pct_idx < len(row_data) else '')
            if pct >= 100:
                _set_cell_bg(tbl.cell(ri, exec_pct_idx), GREEN_100)
            elif pct >= 90:
                _set_cell_bg(tbl.cell(ri, exec_pct_idx), GREEN_HIGH)
            elif pct >= 70:
                _set_cell_bg(tbl.cell(ri, exec_pct_idx), GREEN_MED)
            elif pct >= 50:
                _set_cell_bg(tbl.cell(ri, exec_pct_idx), YELLOW)

        # Color code Test Status
        if status_idx >= 0 and status_idx < len(tbl.columns):
            status = str(row_data.iloc[status_idx] if status_idx < len(row_data) else '').strip()
            if status == 'Test Complete':
                _set_cell_bg(tbl.cell(ri, status_idx), GREEN_100)


def add_table_parts_colored(prs, policy, df, title, logo_path):
    """Add table slides with conditional color coding for Features in Scope."""
    # First add the table normally
    slides_before = len(prs.slides)
    added = add_table_parts(prs, policy, df, title, logo_path)
    # Now color-code each table slide that was just added
    for si in range(slides_before, len(prs.slides)):
        slide = prs.slides[si]
        for shape in slide.shapes:
            if shape.shape_type == 19:  # table
                tbl = shape.table
                # Figure out which chunk of df this table represents
                data_rows = len(tbl.rows) - 1  # minus header
                start_row = (si - slides_before) * policy.rows_per_part
                # Get the corresponding df slice
                chunk = df.iloc[start_row:start_row + data_rows].reset_index(drop=True) if start_row < len(df) else pd.DataFrame()
                if not chunk.empty:
                    _color_code_features_table(tbl, chunk)
    return added


def add_cycle_report_slides(prs, policy, cycle_results, logo_path):
    """Add pie chart screenshot slides (both on one slide)."""
    if not cycle_results: return

    from PIL import Image as PILImage

    prog_img = cycle_results.get('progression_pie_img')
    reg_img = cycle_results.get('regression_pie_img')
    if not prog_img and not reg_img: return

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    SIDE = Inches(0.60); usable_w = prs.slide_width - Inches(1.20)
    tb = s.shapes.add_textbox(SIDE, Inches(0.30), usable_w, Inches(0.30))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = False
    tf.text = 'Test Execution Summary — Progression & Regression (Pie Charts)'
    p = tf.paragraphs[0]; p.font.name = FONT_NAME; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = NAVY
    IMG_TOP = Inches(0.70)
    half_w = int((prs.slide_width - Inches(1.40)) / 2)
    max_h = prs.slide_height - IMG_TOP - Inches(0.70)
    def _add_pie(img_path, left_pos, subtitle):
        if not img_path or not Path(img_path).exists(): return
        with PILImage.open(str(img_path)) as img:
            iw, ih = img.size
        aspect = iw / ih
        if half_w / max_h > aspect: h = max_h; w = int(h * aspect)
        else: w = half_w; h = int(w / aspect)
        x = left_pos + int((half_w - w) / 2)
        s.shapes.add_picture(str(img_path), x, IMG_TOP, w, h)
        sub_top = IMG_TOP + h + Inches(0.05)
        stb = s.shapes.add_textbox(left_pos, sub_top, half_w, Inches(0.25))
        stf = stb.text_frame; stf.clear(); stf.word_wrap = False; stf.text = subtitle
        sp = stf.paragraphs[0]; sp.font.name = FONT_NAME; sp.font.size = Pt(10); sp.font.bold = True; sp.font.color.rgb = NAVY
        sp.alignment = PP_ALIGN.CENTER
    _add_pie(prog_img, Inches(0.50), 'Progression')
    _add_pie(reg_img, Inches(0.50) + half_w + Inches(0.40), 'Regression')
    add_brand_footer_logo(s, prs, logo_path)

    # Features in Scope table (editable, with auto-split)
    if cycle_results.get('features_df') is not None and not cycle_results['features_df'].empty:
        add_table_parts_colored(prs, policy, cycle_results['features_df'],
                       'Features in Scope', logo_path)


# ═══════════════════════════════════════════════════════════════════════════════
# EMBEDDED: slide2_builder.py
# ═══════════════════════════════════════════════════════════════════════════════

_S2_TEMPLATE_PATH = Path(r'C:\Users\P3314665\Downloads\MDA QA INTG_Slide2.pptx')

def _s2_clone_slide(src_prs, src_slide, dst_prs, insert_index):
    layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(layout)
    for ph in list(new_slide.placeholders):
        sp = ph._element; sp.getparent().remove(sp)
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    src_bg = src_slide._element.find(qn('p:bg'))
    if src_bg is not None:
        dst_bg = new_slide._element.find(qn('p:bg'))
        if dst_bg is not None: new_slide._element.remove(dst_bg)
        new_slide._element.insert(0, copy.deepcopy(src_bg))
    sldIdLst = dst_prs.slides._sldIdLst
    sldId_items = list(sldIdLst)
    last = sldId_items[-1]
    sldIdLst.remove(last)
    if insert_index >= len(sldId_items): sldIdLst.append(last)
    else: sldIdLst.insert(insert_index, last)
    return new_slide

def _s2_rescale_to_fit(slide, src_width, dst_width, dst_height):
    if src_width <= dst_width: return
    scale = dst_width / src_width
    for s in slide.shapes:
        s.left = int(s.left * scale)
        s.width = max(int(s.width * scale), 1)
        if s.shape_type == 19:
            tbl = s.table
            total_old = sum(tbl.columns[c].width for c in range(len(tbl.columns)))
            if total_old > 0:
                tbl_scale = s.width / total_old
                for c in range(len(tbl.columns)):
                    tbl.columns[c].width = int(tbl.columns[c].width * tbl_scale)

def _s2_shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name: return s
    return None

def _s2_set_left_panel_text(slide, iteration):
    rect = _s2_shape_by_name(slide, 'Rectangle 2')
    if rect is None: return
    parts = iteration.split('.')
    major, minor = int(parts[0]), int(parts[1])
    next_iter = f'{major + 1}.{minor}'
    tf = rect.text_frame; body = tf._txBody
    lines = [
        (0, 'Offshore team is actively engaged in T-Mobile testing', None),
        (1, f'Integration Progression {next_iter}: Execution', ' - ', 'In Progress'),
        (1, f'Integration Regression {next_iter}: Execution ', '- ', 'In Progress'),
        (1, f'Integration Regression {iteration}: Execution ', '- ', 'Completed'),
        (1, f'Integration Progression {iteration}: Execution', ' - ', 'Completed'),
        (None, '', None),
        (0, 'Automation regression testing update', None),
        (1, 'Yet to have', None),
    ]
    existing_paras = body.findall(qn('a:p'))
    tmpl_l0 = copy.deepcopy(existing_paras[0]) if len(existing_paras) > 0 else None
    tmpl_l1 = copy.deepcopy(existing_paras[1]) if len(existing_paras) > 1 else None
    tmpl_blank = copy.deepcopy(existing_paras[5]) if len(existing_paras) > 5 else None
    for p in existing_paras: body.remove(p)
    def _make_run(text, bold=None, size=1000):
        r = etree.SubElement(etree.Element('dummy'), qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', 'en-US'); rPr.set('sz', str(size)); rPr.set('dirty', '0')
        if bold: rPr.set('b', '1')
        fill = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(fill, qn('a:srgbClr')); clr.set('val', '000000')
        latin = etree.SubElement(rPr, qn('a:latin')); latin.set('typeface', 'Arial')
        cs = etree.SubElement(rPr, qn('a:cs')); cs.set('typeface', 'Arial')
        t = etree.SubElement(r, qn('a:t')); t.text = text
        return r
    for line_def in lines:
        if line_def[0] is None:
            if tmpl_blank is not None:
                p = copy.deepcopy(tmpl_blank)
                for r in p.findall(qn('a:r')): p.remove(r)
            else: p = etree.SubElement(body, qn('a:p'))
            body.append(p); continue
        level = line_def[0]
        tmpl = tmpl_l1 if level == 1 else tmpl_l0
        if tmpl is None: tmpl = tmpl_l0 or etree.Element(qn('a:p'))
        p = copy.deepcopy(tmpl)
        for r in p.findall(qn('a:r')): p.remove(r)
        for epr in p.findall(qn('a:endParaRPr')): p.remove(epr)
        if len(line_def) == 3:
            p.append(_make_run(line_def[1]))
        elif len(line_def) == 4:
            p.append(_make_run(line_def[1]))
            p.append(_make_run(line_def[2], bold=True))
            p.append(_make_run(line_def[3]))
        body.append(p)

def _s2_clear_and_set_cell(tbl, row, col, text, bold=False, size=101600):
    cell = tbl.cell(row, col); cell.text = ''
    tf = cell.text_frame; p = tf.paragraphs[0]; p.clear()
    run = p.add_run(); run.text = str(text)
    run.font.name = 'Arial'; run.font.size = Pt(size // 12700)
    run.font.bold = bold; run.font.color.rgb = RGBColor(0, 0, 0)

def _s2_populate_manual_table(slide, labels_df, iteration, exec_df=None):
    tbl_shape = _s2_shape_by_name(slide, 'Table 1')
    if tbl_shape is None: return
    tbl = tbl_shape.table
    header_map = {}
    for ci in range(len(tbl.columns)):
        header_map[tbl.cell(0, ci).text.strip()] = ci
    parts = iteration.split('.'); major, minor = int(parts[0]), int(parts[1])
    cur_iter = iteration
    col_aliases = {
        'Total': ['TOTAL', 'Total', 'total'],
        'Passed': ['PASS', 'Passed', 'Pass', 'passed'],
        'Failed': ['FAIL', 'Failed', 'Fail', 'failed'],
        'In Progress': ['IN PROGRESS', 'WORK IN PROGRESS', 'In Progress', 'in progress'],
        'No Run': ['NOT EXECUTED', 'No Run', 'No_Run', 'no run'],
        'Not Applicable': ['NOT APPLICABLE', 'Not Applicable', 'not applicable'],
        'Blocked': ['BLOCKED', 'Blocked', 'blocked'],
    }
    def _find_col(df, aliases):
        for a in aliases:
            if a in df.columns: return a
        return None
    def _sum_col(df, aliases):
        c = _find_col(df, aliases)
        if c is None: return 0
        try: return int(df[c].sum())
        except Exception: return 0
    exec_totals = {}
    if exec_df is not None and not exec_df.empty:
        assignee_col = None
        for c in exec_df.columns:
            if c.strip().upper() in ('EXECUTION ASSIGNEE', 'ASSIGNEE'): assignee_col = c; break
        total_row = None
        if assignee_col:
            mask = exec_df[assignee_col].astype(str).str.strip().str.upper() == 'TOTAL'
            if mask.any(): total_row = exec_df[mask].iloc[0]
        if total_row is not None:
            for hdr, aliases in col_aliases.items():
                col = _find_col(exec_df, aliases)
                if col:
                    try: exec_totals[hdr] = int(total_row[col])
                    except Exception: exec_totals[hdr] = 0
                else: exec_totals[hdr] = 0
        else:
            prog_pattern = f'PROGRESSION.*{re.escape(cur_iter)}'
            exec_sub = exec_df
            if 'Release' in exec_df.columns:
                mask = exec_df['Release'].astype(str).str.contains(prog_pattern, case=False, na=False)
                if mask.any(): exec_sub = exec_df[mask]
            for hdr, aliases in col_aliases.items():
                exec_totals[hdr] = _sum_col(exec_sub, aliases)
    row_defs = [
        (f'Integration progression {cur_iter}', True),
        (f'Integration Regression {cur_iter}', False),
        (f'Prod defects {major}.X', False),
    ]
    for ri, (label, populate) in enumerate(row_defs, start=1):
        if ri >= len(tbl.rows): break
        _s2_clear_and_set_cell(tbl, ri, 0, label)
        if populate and exec_totals:
            for hdr, aliases in col_aliases.items():
                ci = header_map.get(hdr)
                if ci is None: continue
                val = exec_totals.get(hdr, 0)
                _s2_clear_and_set_cell(tbl, ri, ci, val if val else '-', bold=(hdr == 'Total'))
        else:
            for ci in range(1, len(tbl.columns)):
                _s2_clear_and_set_cell(tbl, ri, ci, '-')

_S2_PRIORITY_MAP = {
    'P0': ['blocker', 'emergency', 'blocker/emergency', 'p0'],
    'P1': ['critical', 'high', 'critical/high', 'p1'],
    'P2': ['major', 'medium', 'major/medium', 'p2'],
    'P3': ['minor', 'low', 'minor/low', 'p3'],
    'P4': ['trivial', 'p4'],
}

def _s2_build_priority_totals(defects_df):
    result = {p: 0 for p in ['P0', 'P1', 'P2', 'P3', 'P4']}
    if defects_df is None or defects_df.empty: return result
    cols_upper = {c.upper().strip(): c for c in defects_df.columns}
    if 'P0' in cols_upper or 'P1' in cols_upper:
        for p in result:
            real = cols_upper.get(p)
            if real:
                try: result[p] = int(defects_df[real].sum())
                except Exception: pass
        return result
    _COL_TO_P = {
        'blocker/emergency': 'P0', 'blocker': 'P0', 'emergency': 'P0',
        'critical/high': 'P1', 'critical': 'P1', 'high': 'P1',
        'major/medium': 'P2', 'major': 'P2', 'medium': 'P2',
        'minor/low': 'P3', 'minor': 'P3', 'low': 'P3',
        'trivial': 'P4',
    }
    matched_any = False
    for col in defects_df.columns:
        col_lower = col.strip().lower()
        p_key = _COL_TO_P.get(col_lower)
        if p_key:
            matched_any = True
            try: result[p_key] += int(defects_df[col].sum())
            except Exception: pass
    if matched_any: return result
    pri_col = cols_upper.get('PRIORITY') or cols_upper.get('PRIORITY NAME')
    cnt_col = cols_upper.get('COUNT') or cols_upper.get('TOTAL') or cols_upper.get('DEFECT COUNT')
    if pri_col and cnt_col:
        for _, row in defects_df.iterrows():
            pri_val = str(row[pri_col]).strip().lower()
            try: cnt = int(row[cnt_col])
            except Exception: cnt = 0
            for p_key, aliases in _S2_PRIORITY_MAP.items():
                if any(a in pri_val for a in aliases): result[p_key] += cnt; break
    elif pri_col:
        for _, row in defects_df.iterrows():
            pri_val = str(row[pri_col]).strip().lower()
            for p_key, aliases in _S2_PRIORITY_MAP.items():
                if any(a in pri_val for a in aliases): result[p_key] += 1; break
    return result

def _s2_populate_defect_table(slide, defects_df, iteration):
    tbl_shape = _s2_shape_by_name(slide, 'Table 4')
    if tbl_shape is None: return
    tbl = tbl_shape.table
    header_map = {}
    for ci in range(len(tbl.columns)):
        header_map[tbl.cell(0, ci).text.strip()] = ci
    parts = iteration.split('.'); major = int(parts[0]); cur_iter = iteration
    row_labels = [
        (f'Integration progression {cur_iter}', True),
        (f'Integration  Regression {cur_iter}', False),
        (f'Automation Regression {major} X', False),
        ('Smoke testing', False),
    ]
    p_totals = _s2_build_priority_totals(defects_df)
    for ri, (label, populate) in enumerate(row_labels, start=1):
        if ri >= len(tbl.rows): break
        _s2_clear_and_set_cell(tbl, ri, 0, label)
        if populate:
            grand = 0
            for p_key in ['P0', 'P1', 'P2', 'P3', 'P4']:
                ci = header_map.get(p_key); val = p_totals.get(p_key, 0)
                if ci is not None: _s2_clear_and_set_cell(tbl, ri, ci, val if val else '-')
                grand += val
            ci_total = header_map.get('Total')
            if ci_total is not None:
                _s2_clear_and_set_cell(tbl, ri, ci_total, grand if grand else '-', bold=True)
        else:
            for ci in range(1, len(tbl.columns)):
                _s2_clear_and_set_cell(tbl, ri, ci, '-')

def _s2_update_date_textbox(slide):
    tb = _s2_shape_by_name(slide, 'TextBox 23')
    if tb is None: return
    now = datetime.now()
    mon = now - timedelta(days=now.weekday()); fri = mon + timedelta(days=4)
    def _ordinal(d):
        day = d.day
        suffix = 'th' if 11 <= day <= 13 else {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')
        return f'{day}{suffix}'
    date_text = f"Manual ({_ordinal(mon)} {mon.strftime('%b')} to {_ordinal(fri)} {fri.strftime('%b')})"
    tf = tb.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    p = tf.paragraphs[0]; p.clear()
    run = p.add_run(); run.text = date_text
    run.font.name = 'Arial'; run.font.size = Pt(9); run.font.bold = True

def add_slide2_delivery_updates(dst_prs, iteration, labels_df=None, defects_df=None,
                                exec_df=None, template_path=_S2_TEMPLATE_PATH, insert_index=1):
    if not template_path.exists():
        print(f'[WARN] Slide2 template not found: {template_path}'); return
    src_prs = Presentation(str(template_path))
    src_slide = src_prs.slides[0]
    new_slide = _s2_clone_slide(src_prs, src_slide, dst_prs, insert_index)
    _s2_rescale_to_fit(new_slide, src_prs.slide_width, dst_prs.slide_width, dst_prs.slide_height)
    _s2_set_left_panel_text(new_slide, iteration)
    _s2_update_date_textbox(new_slide)
    if labels_df is not None or exec_df is not None:
        _s2_populate_manual_table(new_slide, labels_df, iteration, exec_df=exec_df)
    if defects_df is not None:
        _s2_populate_defect_table(new_slide, defects_df, iteration)
    print(f'[OK] Slide 2 (Delivery Updates) inserted at position {insert_index + 1}')
    return new_slide

st.set_page_config(page_title='MDA Jira Dashboard', page_icon=':bar_chart:', layout='wide')

# ---------------- CSS (anchors + theme) ----------------
st.markdown(
    '''<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&display=swap');

    html, body { margin:0 !important; padding:0 !important; scroll-behavior: smooth; }

    /* Animations */
    @keyframes fadeSlideIn { from { opacity:0; transform:translateY(18px); } to { opacity:1; transform:translateY(0); } }
    @keyframes pulse { 0%,100% { opacity:1; } 50% { opacity:0.6; } }
    @keyframes shimmer { 0% { background-position: -200% 0; } 100% { background-position: 200% 0; } }
    @keyframes float { 0%,100% { transform: translateY(0); } 50% { transform: translateY(-6px); } }
    @keyframes glow { 0%,100% { box-shadow: 0 0 20px rgba(124,58,237,0.3); } 50% { box-shadow: 0 0 35px rgba(59,130,246,0.5); } }

    /* App background - dark charcoal */
    .stApp {
        background: #0e1117 !important;
        color: #f0f2f6;
        font-family: "Inter", -apple-system, BlinkMacSystemFont, sans-serif !important;
    }

    /* Hide Streamlit chrome */
    [data-testid="stToolbar"], header[data-testid="stHeader"], #MainMenu, footer, header[tabindex="-1"]
    { display:none !important; visibility:hidden !important; height:0 !important; margin:0 !important; padding:0 !important; }
    .block-container { padding-top: 0 !important; padding-left: 130px !important; padding-right: 24px !important; max-width: 100% !important; overflow: visible; }

    /* Left Rail - Neon glass */
    .left-rail {
        position: fixed; left: 16px; top: 18px; bottom: 18px; width: 72px; z-index: 10;
        background: linear-gradient(180deg, #0a2a2a 0%, #0c3535 100%);
        border-radius: 24px;
        border: 1px solid rgba(6,182,212,0.3);
        box-shadow: 0 8px 32px rgba(0,0,0,0.5), 0 0 40px rgba(6,182,212,0.1);
        backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
        display: flex; flex-direction: column; align-items: center; padding: 18px 12px; gap: 14px;
        animation: fadeSlideIn 0.4s ease-out;
    }
    .rail-spacer { flex: 1; }
    .rail-link { text-decoration: none; }
    .rail-item {
        width: 48px; height: 48px; display: flex; align-items: center; justify-content: center;
        background: rgba(255,255,255,0.12); color: #FFFFFF;
        border: 1px solid rgba(255,255,255,0.25); border-radius: 14px;
        backdrop-filter: blur(8px); box-shadow: inset 0 1px 0 rgba(255,255,255,0.15);
        transition: all 0.25s cubic-bezier(0.4,0,0.2,1);
    }
    .rail-item:hover {
        transform: scale(1.15) translateY(-2px);
        background: rgba(6,182,212,0.3);
        box-shadow: 0 8px 25px rgba(6,182,212,0.4);
    }
    .rail-item svg { filter: drop-shadow(0 2px 4px rgba(0,0,0,0.2)); }

    /* Content wrapper */
    .content-wrap { margin-left: 0; padding: 18px 20px 28px 20px; position: relative; z-index: 1; }

    /* Top Banner */
    .top-banner {
        display: flex; align-items: center; justify-content: space-between; flex-wrap: wrap;
        background: #1a1d24;
        padding: 18px 24px;
        border: 1px solid rgba(255,255,255,0.10);
        border-radius: 20px;
        box-shadow: 0 4px 16px rgba(0,0,0,0.4);
        animation: fadeSlideIn 0.5s ease-out;
    }
    .title {
        font-weight: 800; font-size: 26px; letter-spacing: -0.5px;
        background: linear-gradient(135deg, #ffffff 0%, #22d3ee 100%);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        background-clip: text;
    }
    .sub { color: #b8c0cc; font-size: 13px; font-weight: 500; margin-top: 2px; }

    /* Pills */
    .pill {
        display: inline-flex; align-items: center; gap: 6px;
        background: rgba(6,182,212,0.18);
        border: 1px solid rgba(6,182,212,0.35);
        padding: 8px 14px; border-radius: 12px;
        font-size: 12px; font-weight: 600; color: #ffffff;
        transition: all 0.2s ease;
    }
    .pill:hover { background: rgba(6,182,212,0.3); transform: translateY(-1px); }

    /* Stats Grid */
    .stats-grid { display: grid; grid-template-columns: repeat(4, minmax(0,1fr)); gap: 16px; margin-top: 16px; }
    .stat {
        border-radius: 18px; padding: 18px; position: relative; overflow: hidden;
        background: #1a1d24;
        border: 1px solid rgba(255,255,255,0.10);
        box-shadow: 0 4px 16px rgba(0,0,0,0.3);
        animation: fadeSlideIn 0.6s ease-out;
        transition: all 0.3s cubic-bezier(0.4,0,0.2,1);
    }
    .stat::before {
        content: ""; position: absolute; top: 0; left: 0; right: 0; height: 3px;
        background: linear-gradient(90deg, var(--accent-start), var(--accent-end));
        opacity: 0.8;
    }
    .stat:hover { transform: translateY(-4px) scale(1.02); box-shadow: 0 8px 24px rgba(0,0,0,0.4); }
    .stat-purple { --accent-start: #06b6d4; --accent-end: #0891b2; }
    .stat-blue { --accent-start: #0891b2; --accent-end: #0e7490; }
    .stat-rose { --accent-start: #22d3ee; --accent-end: #06b6d4; }
    .stat-amber { --accent-start: #67e8f9; --accent-end: #22d3ee; }
    .stat .label { font-size: 12px; color: #c8d1dc; font-weight: 600; text-transform: uppercase; letter-spacing: 0.5px; }
    .stat .value { font-size: 32px; font-weight: 800; color: #FFFFFF; margin-top: 4px; text-shadow: 0 0 12px rgba(6,182,212,0.25); }

    /* Cards */
    .card {
        background: #1a1d24;
        border: 1px solid rgba(255,255,255,0.10);
        border-radius: 20px; padding: 20px 22px;
        box-shadow: 0 4px 16px rgba(0,0,0,0.3);
        animation: fadeSlideIn 0.65s ease-out;
    }
    .section-grid { display: grid; grid-template-columns: 2.2fr 1.2fr; gap: 16px; margin-top: 16px; }

    /* CLI */
    .cli-header {
        font-weight: 700; color: #ffffff; letter-spacing: 0.3px;
        background: rgba(6,182,212,0.15);
        border: 1px solid rgba(6,182,212,0.3);
        border-radius: 12px; padding: 10px 14px; margin: 0 0 10px 0;
        display: flex; align-items: center; gap: 8px;
    }
    .cli-header::before { content: "●"; color: #22d3ee; font-size: 10px; animation: pulse 2s infinite; }
    .cli-box {
        background: #141720;
        color: #c8d1dc;
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 14px; padding: 14px 16px;
        height: 380px; max-height: 380px; overflow: auto;
        box-shadow: inset 0 2px 8px rgba(0,0,0,0.4);
    }
    .cli-box pre { margin: 0; font-family: "JetBrains Mono", ui-monospace, SFMono-Regular, Menlo, Consolas, monospace; font-size: 12px; line-height: 1.5; }
    .cli-box::-webkit-scrollbar { width: 6px; }
    .cli-box::-webkit-scrollbar-track { background: rgba(255,255,255,0.05); border-radius: 3px; }
    .cli-box::-webkit-scrollbar-thumb { background: rgba(6,182,212,0.4); border-radius: 3px; }

    /* Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #0e7490 0%, #22d3ee 100%) !important;
        color: #FFFFFF !important; border: 0 !important;
        padding: 0.85rem 1.5rem !important; border-radius: 14px !important;
        font-weight: 700 !important; font-size: 14px !important; letter-spacing: 0.3px;
        box-shadow: 0 6px 20px rgba(6,182,212,0.3) !important;
        transition: all 0.3s cubic-bezier(0.4,0,0.2,1) !important;
    }
    .stButton > button:hover {
        transform: translateY(-3px) scale(1.02) !important;
        box-shadow: 0 10px 30px rgba(6,182,212,0.45) !important;
    }
    .stButton > button:active { transform: translateY(-1px) scale(0.98) !important; }

    /* Form elements */
    [data-testid="stCheckbox"] label { color: #ffffff !important; font-weight: 600 !important; }
    [data-testid="stCheckbox"] span[data-testid="stCheckboxLabel"] { color: #ffffff !important; }
    .stTextInput label, .stCaption { color: #c8d1dc !important; font-weight: 600 !important; }
    .stTextInput input {
        color: #000000 !important;
        background: rgba(255,255,255,0.9) !important;
        border: 1px solid rgba(255,255,255,0.15) !important;
        border-radius: 12px !important;
        font-weight: 500 !important;
        transition: all 0.2s ease !important;
    }
    .stTextInput input:focus {
        border-color: rgba(6,182,212,0.5) !important;
        box-shadow: 0 0 15px rgba(6,182,212,0.15) !important;
    }

    /* Summary card */
    .summary {
        background: rgba(6,182,212,0.1);
        border: 1px solid rgba(6,182,212,0.25);
        border-radius: 18px; padding: 16px 18px;
        box-shadow: 0 4px 16px rgba(0,0,0,0.2);
        animation: fadeSlideIn 0.7s ease-out;
    }
    .summary h4 { margin: 0 0 8px 0; color: #22d3ee; font-weight: 700; }
    .summary ul { margin: 0; padding-left: 20px; color: #c8d1dc; }
    .summary li { margin: 4px 0; }

    /* Section titles */
    .section-title {
        font-size: 14px; font-weight: 700; text-transform: uppercase; letter-spacing: 1px;
        color: #c8d1dc; margin-bottom: 12px;
        display: flex; align-items: center; gap: 8px;
    }
    .section-title::after { content: ""; flex: 1; height: 1px; background: linear-gradient(90deg, rgba(255,255,255,0.1), transparent); }

    /* Download buttons */
    .stDownloadButton > button {
        background: rgba(255,255,255,0.06) !important;
        border: 1px solid rgba(255,255,255,0.1) !important;
        color: #ffffff !important;
        border-radius: 10px !important;
        font-weight: 600 !important;
        transition: all 0.2s ease !important;
    }
    .stDownloadButton > button:hover {
        background: rgba(6,182,212,0.15) !important;
        border-color: rgba(6,182,212,0.4) !important;
        transform: translateY(-2px) !important;
    }

    /* Checkbox styling - box */
    [data-testid="stCheckbox"] > label > div:first-child {
        background: rgba(255,255,255,0.1) !important;
        border: 1px solid rgba(255,255,255,0.2) !important;
        border-radius: 6px !important;
    }
    [data-testid="stCheckbox"] > label > div:first-child[data-checked="true"] {
        background: linear-gradient(135deg, #0e7490, #22d3ee) !important;
        border-color: transparent !important;
    }
    /* Checkbox labels - force light text */
    [data-testid="stCheckbox"] label,
    [data-testid="stCheckbox"] label span,
    [data-testid="stCheckbox"] label p,
    [data-testid="stCheckbox"] p,
    [data-testid="stCheckbox"] span {
        color: #ffffff !important;
        font-weight: 500 !important;
    }
    </style>''',
    unsafe_allow_html=True,
)

# ---------------- State ----------------
ss = st.session_state
for k, v in [('reset_mods', False), ('reset_all', False), ('mod1', False), ('mod2', False), ('mod3', False), ('last_summary', None), ('cli_lines', []), ('cli_error', None), ('select_all', False)]:
    if k not in ss: ss[k] = v
if ss.get('reset_mods') or ss.get('reset_all'):
    ss['mod1'] = ss['mod2'] = ss['mod3'] = False
    ss['select_all'] = False
    if ss.get('reset_all'):
        ss['cli_lines'] = []
        ss['cli_error'] = None
        ss['last_summary'] = None
    ss['reset_mods'] = False
    ss['reset_all'] = False

# Pre-read iteration so banner and pills use the freshest value
if 'iteration' not in ss:
    ss['iteration'] = '51.2'

# ---------------- Left rail with anchors ----------------
st.markdown(
    """
<div class='left-rail'>
  <a class='rail-link' href='#dashboard'><div class='rail-item' title='Dashboard'><svg width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M3 12l2-2m0 0l7-7 7 7M5 10v10a1 1 0 001 1h3m10-11l2 2m-2-2v10a1 1 0 01-1 1h-3m-6 0h6"/></svg></div></a>
  <a class='rail-link' href='#modules'><div class='rail-item' title='Modules'><svg width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><rect x="3" y="3" width="7" height="7" rx="1"/><rect x="14" y="3" width="7" height="7" rx="1"/><rect x="3" y="14" width="7" height="7" rx="1"/><rect x="14" y="14" width="7" height="7" rx="1"/></svg></div></a>
  <a class='rail-link' href='#cli'><div class='rail-item' title='CLI'><svg width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><rect x="2" y="4" width="20" height="16" rx="2"/><path d="M6 9l4 3-4 3M12 16h6"/></svg></div></a>
  <a class='rail-link' href='#outputs'><div class='rail-item' title='Outputs'><svg width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M20 7l-8-4-8 4m16 0l-8 4m8-4v10l-8 4m0-10L4 7m8 4v10M4 7v10l8 4"/></svg></div></a>
  <div class='rail-spacer'></div>
  <a class='rail-link' href='#dashboard'><div class='rail-item' title='Profile'><svg width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><circle cx="12" cy="8" r="4"/><path d="M4 20c0-4 4-6 8-6s8 2 8 6"/></svg></div></a>
</div>
""",
    unsafe_allow_html=True,
)

# ---------------- Content ----------------
st.markdown("<div class='content-wrap'>", unsafe_allow_html=True)

# Anchored Dashboard banner
st.markdown("<div id='dashboard'></div>", unsafe_allow_html=True)

from datetime import datetime as _dt
mon = (_dt.now() - timedelta(days=_dt.now().weekday())).strftime('%d %b %Y')
fri = (_dt.now() + timedelta(days=(4 - _dt.now().weekday()))).strftime('%d %b %Y')
st.markdown(
    f"""
<div class='top-banner'>
  <div>
    <div class='title'>MDA Jira Dashboard</div>
    <div class='sub'>Jira QA Status &mdash; {mon} to {fri}</div>
  </div>
  <div style='display:flex; gap:10px; align-items:center;'>
    <div class='pill'>
      <svg width="14" height="14" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M12 2L2 7l10 5 10-5-10-5zM2 17l10 5 10-5M2 12l10 5 10-5"/></svg>
      Iteration {escape(str(ss.get('iteration','51.2')))}
    </div>
    <div class='pill'>
      <svg width="14" height="14" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M9 19v-6a2 2 0 00-2-2H5a2 2 0 00-2 2v6m6 0h6m-6 0V9a2 2 0 012-2h2a2 2 0 012 2v10m6 0v-4a2 2 0 00-2-2h-2a2 2 0 00-2 2v4"/></svg>
      Reports
    </div>
    <div class='pill'>
      <svg width="14" height="14" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M9 12l2 2 4-4m6 2a9 9 0 11-18 0 9 9 0 0118 0z"/></svg>
      QA
    </div>
  </div>
</div>
""",
    unsafe_allow_html=True,
)

# Stats row
completed_items = len(ss['last_summary']['completed']) if ss.get('last_summary') and ss['last_summary'].get('completed') else 0
artifacts_count = sum(1 for k in ['labels_pivot','defects_pivot','exec_excel','ppt'] if ss.get('last_summary') and ss['last_summary'].get(k))
last_ts = (ss['last_summary']['ts'] if ss.get('last_summary') and ss['last_summary'].get('ts') else '-')
selected_now = int(ss['mod1']) + int(ss['mod2']) + int(ss['mod3'])

st.markdown(
    f"""
<div class='stats-grid'>
  <div class='stat stat-purple'>
    <div class='label'><svg width="14" height="14" fill="none" stroke="#06b6d4" stroke-width="2" viewBox="0 0 24 24" style="vertical-align:-2px;margin-right:4px;"><rect x="3" y="3" width="7" height="7" rx="1"/><rect x="14" y="3" width="7" height="7" rx="1"/><rect x="3" y="14" width="7" height="7" rx="1"/><rect x="14" y="14" width="7" height="7" rx="1"/></svg>Selected Modules</div>
    <div class='value'>{selected_now}</div>
  </div>
  <div class='stat stat-blue'>
    <div class='label'><svg width="14" height="14" fill="none" stroke="#0891b2" stroke-width="2" viewBox="0 0 24 24" style="vertical-align:-2px;margin-right:4px;"><path d="M20 7l-8-4-8 4m16 0l-8 4m8-4v10l-8 4m0-10L4 7m8 4v10M4 7v10l8 4"/></svg>Artifacts (last run)</div>
    <div class='value'>{artifacts_count}</div>
  </div>
  <div class='stat stat-rose'>
    <div class='label'><svg width="14" height="14" fill="none" stroke="#22d3ee" stroke-width="2" viewBox="0 0 24 24" style="vertical-align:-2px;margin-right:4px;"><path d="M9 12l2 2 4-4m6 2a9 9 0 11-18 0 9 9 0 0118 0z"/></svg>Completed Items (last run)</div>
    <div class='value'>{completed_items}</div>
  </div>
  <div class='stat stat-amber'>
    <div class='label'><svg width="14" height="14" fill="none" stroke="#67e8f9" stroke-width="2" viewBox="0 0 24 24" style="vertical-align:-2px;margin-right:4px;"><circle cx="12" cy="12" r="10"/><path d="M12 6v6l4 2"/></svg>Last Run</div>
    <div class='value'>{escape(str(last_ts))}</div>
  </div>
</div>
""",
    unsafe_allow_html=True,
)

# Main split
st.markdown("<div class='section-grid'>", unsafe_allow_html=True)

# LEFT COLUMN (Modules)
st.markdown("<div id='modules'></div>", unsafe_allow_html=True)
left_col = st.container()
with left_col:
    st.markdown("<div class='section-title'>Automation Modules</div>", unsafe_allow_html=True)
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    c1, c2 = st.columns([1,1])
    with c1:
        iteration = st.text_input('Iteration (e.g., 51.2)', value=ss.get('iteration','51.2'))
        iteration = iteration.strip() or '51.2'
        ss['iteration'] = iteration  # updates pill dynamically on rerun
    with c2:
        headed = st.checkbox('Show Browser (headed)', value=False)
        os.environ['HEADLESS'] = 'false' if headed else 'true'

    st.caption('Select Modules to Run')
    def _on_select_all():
        val = ss['select_all']
        ss['mod1'] = ss['mod2'] = ss['mod3'] = val
    st.checkbox('✅ Select All Modules', key='select_all', on_change=_on_select_all)
    mcol1, mcol2 = st.columns([1,1])
    with mcol1:
        m1 = st.checkbox("📥 1) QMetry Extract + Pivot (Labels)", key='mod1')
        m2 = st.checkbox("🐞 2) Jira Defect Extract + Pivot (Defects)", key='mod2')
    with mcol2:
        m3 = st.checkbox("📈 3) Jira -> Test Execution Report + Excel", key='mod3')
        st.caption('4) PPT will auto-run only when 1, 2 and 3 complete successfully')

    btn_col1, btn_col2, btn_col3, btn_col4 = st.columns([1.2, 0.6, 0.6, 0.6])
    with btn_col1:
        run_btn = st.button('Run Selected Modules', type='primary')
    with btn_col2:
        clear_btn = st.button('🧹 Clear', key='btn_clear')
    with btn_col3:
        reload_btn = st.button('🔄 Reload', key='btn_reload')
    with btn_col4:
        exit_btn = st.button('🚪 Exit', key='btn_exit')
    st.markdown("</div>", unsafe_allow_html=True)

    # Completion Summary (persistent)
    if ss.get('last_summary'):
        s = ss['last_summary']
        st.markdown("<div class='summary'>" +
                    f"<h4>✅ Completed — Iteration {escape(str(s['iteration']))} <small style='float:right'>{escape(s['ts'])}</small></h4>" +
                    "<ul>" + ''.join(f"<li>{escape(item)}</li>" for item in (s['completed'] or ['(No artifacts)'])) + "</ul>" +
                    "</div>", unsafe_allow_html=True)
        sc1, sc2 = st.columns(2)
        from pathlib import Path as _P
        with sc1:
            p = s.get('labels_pivot');
            if p and _P(p).exists():
                st.download_button('⬇️ Download Labels Pivot', data=_P(p).read_bytes(), file_name=_P(p).name, key='sum_dl_labels')
            p = s.get('defects_pivot');
            if p and _P(p).exists():
                st.download_button('⬇️ Download Defects Pivot', data=_P(p).read_bytes(), file_name=_P(p).name, key='sum_dl_defects')
        with sc2:
            p = s.get('exec_excel');
            if p and _P(p).exists():
                st.download_button('⬇️ Download Test Execution Excel', data=_P(p).read_bytes(), file_name=_P(p).name, key='sum_dl_exec')
            p = s.get('ppt');
            if p and _P(p).exists():
                st.download_button('⬇️ Download Weekly PPT', data=_P(p).read_bytes(), file_name=_P(p).name, key='sum_dl_ppt')
        if st.button('Dismiss Summary', key='btn_dismiss_sum'):
            ss['last_summary'] = None
            st.rerun()

# Handle Clear button — uncheck all, clear summary + CLI
if clear_btn:
    ss['mod1'] = ss['mod2'] = ss['mod3'] = False
    ss['select_all'] = False
    ss['cli_lines'] = []
    ss['cli_error'] = None
    ss['last_summary'] = None
    st.rerun()

# Handle Reload button — flush caches and rerun
if reload_btn:
    st.cache_resource.clear()
    st.cache_data.clear()
    st.toast('Modules reloaded!')
    st.rerun()

# Handle Exit button — stop the Streamlit script
if exit_btn:
    ss['mod1'] = ss['mod2'] = ss['mod3'] = False
    ss['select_all'] = False
    ss['cli_lines'] = []
    ss['cli_error'] = None
    ss['last_summary'] = None
    st.toast('Exiting dashboard...')
    st.stop()

# RIGHT COLUMN (CLI)
st.markdown("<div id='cli'></div>", unsafe_allow_html=True)
right_col = st.container()
with right_col:
    st.markdown("<div class='section-title'>Live Activity</div>", unsafe_allow_html=True)
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.markdown("<div class='cli-header'>Terminal Output</div>", unsafe_allow_html=True)
    cli_head = st.empty()
    cli_log = st.empty()
    tools = st.empty()
    # Show persisted error from previous run
    if ss.get('cli_error'):
        cli_head.markdown("<div class='cli-header' style='background:linear-gradient(135deg,rgba(239,68,68,0.25),rgba(220,38,38,0.15));border-color:rgba(239,68,68,0.4);color:#FCA5A5;'>❌ %s</div>" % escape(ss['cli_error']), unsafe_allow_html=True)
    # Show persisted log lines
    if ss.get('cli_lines'):
        view = ''.join(reversed(''.join(ss['cli_lines']).splitlines(True)[-1200:]))
        cli_log.markdown("<div class='cli-box'><pre>%s</pre></div>" % escape(view), unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True)  # end section-grid

# Outputs anchor
st.markdown("<div id='outputs'></div>", unsafe_allow_html=True)
outputs_container = st.container()

# ---------------- Helpers ----------------
import re as _re

def validate_iteration(value: str) -> bool:
    return bool(_re.match(r'^\d+\.\d+$', value.strip()))

class LiveLog:
    def __init__(self, head_ph, log_ph, tools_ph):
        self.head = head_ph
        self.log = log_ph
        self.tools = tools_ph
        self.lines = []
        self._ver = 0
    def set(self, text: str):
        self.head.markdown("<div class='cli-header'>>> %s</div>" % escape(text), unsafe_allow_html=True)
    def write(self, s):
        parts = s.splitlines(True)
        if not parts:
            return
        self.lines.extend(parts)
        ss['cli_lines'] = list(self.lines)  # persist for rerun
        view = ''.join(reversed(''.join(self.lines).splitlines(True)[-1200:]))
        self.log.markdown("<div class='cli-box'><pre>%s</pre></div>" % escape(view), unsafe_allow_html=True)
        self._ver += 1
        self.tools.download_button('Download logs.txt', data=''.join(self.lines).encode('utf-8'), file_name='run_logs.txt', key='dl_logs_%d' % self._ver)
    def flush(self):
        pass

# ---------------- Playwright helpers ----------------

def retry_module(func, max_retries=3, delay=5, label='module'):
    """Retry a function up to max_retries times. Retries on exception OR None result."""
    last_err = None
    for attempt in range(1, max_retries + 1):
        try:
            result = func()
            if result is not None:
                return result
            print(f'[RETRY] {label} attempt {attempt}/{max_retries} returned None')
        except Exception as e:
            last_err = e
            print(f'[RETRY] {label} attempt {attempt}/{max_retries} failed: {e}')
        if attempt < max_retries:
            print(f'[RETRY] Waiting {delay}s before retry...')
            time.sleep(delay)
    print(f'[FAIL] {label} failed after {max_retries} attempts')
    return None

def launch_browser_with_fallback(p, headed: bool, logger):
    try:
        return p.chromium.launch(headless=not headed)
    except Exception as e:
        logger.write("[WARN] Chromium launch failed: %s, trying Edge...\n" % e)
        return p.chromium.launch(headless=not headed, channel='msedge')

# ---------------- Execute ----------------
results = {}

if run_btn:
    if not any([m1, m2, m3]):
        st.warning('Please select at least one module to run.')
    elif not validate_iteration(iteration):
        st.error('Iteration must be like 51.2 or 50.3.')
    else:
        logger = LiveLog(cli_head, cli_log, tools)
        logger.set('Initializing...')
        ss['cli_error'] = None  # clear previous error
        ss['cli_lines'] = []    # clear previous log
        # Auto-scroll to CLI window
        st.markdown("<script>document.getElementById('cli').scrollIntoView({behavior:'smooth'});</script>", unsafe_allow_html=True)
        import streamlit.components.v1 as components
        components.html("<script>parent.document.getElementById('cli').scrollIntoView({behavior:'smooth'});</script>", height=0)
        st.success('Starting automation for iteration %s ...' % iteration)
        run_success = False
        _run_t0 = time.time()

        from contextlib import redirect_stdout as _redir
        with _redir(logger):
            try:
                if m1 or m2:
                    with sync_playwright() as p:
                        logger.set('Launching browser ...')
                        browser = launch_browser_with_fallback(p, headed, logger)
                        integration = 'INTEGRATION_PROGRESSION_%s' % iteration
                        if m1:
                            logger.set('\U0001F4E5 QMetry Extract + Pivot (Labels) ...')
                            print('[STEP] Running QMetry Extract + Pivot...')
                            labels_pivot = retry_module(
                                lambda: run_qmetry_automation(browser, iteration, integration),
                                max_retries=3, delay=5, label='QMetry Labels')
                            if labels_pivot:
                                results['labels_pivot'] = Path(str(labels_pivot))
                                print('[OK] Labels pivot => %s' % results['labels_pivot'])
                            else:
                                results['labels_pivot'] = None
                                print('[WARN] QMetry returned None')
                        if m2:
                            logger.set('\U0001F41E Jira Defect Extract + Pivot (Defects) ...')
                            print('[STEP] Running Jira Defect Extract + Pivot...')
                            defects_pivot = retry_module(
                                lambda: run_defect_automation(browser, iteration, integration),
                                max_retries=3, delay=5, label='Jira Defects')
                            if defects_pivot:
                                results['defects_pivot'] = Path(str(defects_pivot))
                                print('[OK] Defects pivot => %s' % results['defects_pivot'])
                            else:
                                results['defects_pivot'] = None
                                print('[WARN] Defect returned None')
                        browser.close()

                if m3:
                    logger.set('\U0001F4C8 Jira -> Test Execution Report + Excel ...')
                    print('[STEP] Running Jira -> Test Execution Report (embedded)...')
                    _qtr_page = None; _qtr_context = None; _qtr_browser = None; _qtr_pw = None
                    _m3_attempts = 0
                    while _m3_attempts < 3:
                        _m3_attempts += 1
                        try:
                            if m1 and m2:
                                ret = run_qmetry_test_report(iteration, keep_browser=True)
                                excel_path, label_excel_path, _qtr_page, _qtr_context, _qtr_browser, _qtr_pw = ret
                            else:
                                excel_path, label_excel_path = run_qmetry_test_report(iteration, keep_browser=False)
                            if excel_path:
                                results['exec_excel'] = Path(str(excel_path))
                                print('[OK] Generated Excel => %s' % results['exec_excel'])
                            else:
                                results['exec_excel'] = None
                            if label_excel_path:
                                results['label_excel'] = Path(str(label_excel_path))
                                print('[OK] Features in Scope => %s' % results['label_excel'])
                            else:
                                results['label_excel'] = None
                            break  # success
                        except Exception as e3:
                            print('[RETRY] Module 3 attempt %d/3 failed: %s' % (_m3_attempts, e3))
                            # Clean up failed browser before retry
                            if _qtr_context:
                                try: _qtr_context.close()
                                except Exception: pass
                            if _qtr_browser:
                                try: _qtr_browser.close()
                                except Exception: pass
                            if _qtr_pw:
                                try: _qtr_pw.stop()
                                except Exception: pass
                            _qtr_page = None; _qtr_context = None; _qtr_browser = None; _qtr_pw = None
                            if _m3_attempts < 3:
                                print('[RETRY] Waiting 5s before retry...')
                                time.sleep(5)
                            else:
                                print('[FAIL] Module 3 failed after 3 attempts')
                                results['exec_excel'] = None
                                results['label_excel'] = None

                # V3.0: Auto-run QMetry Cycle Reports
                if m1 and m2 and m3:
                    logger.set('\U0001F4CA QMetry Cycle Pie Charts ...')
                    print('[STEP] Running QMetry Cycle Reports...')
                    try:
                        cycle_results = run_qmetry_cycle_reports(iteration, page=_qtr_page)
                        results['cycle_reports'] = cycle_results
                        print('[OK] Cycle reports captured')
                    except Exception as ce:
                        print('[WARN] Cycle reports failed: %s' % ce)
                        results['cycle_reports'] = None
                    finally:
                        if _qtr_context:
                            try: _qtr_context.close()
                            except Exception: pass
                        if _qtr_browser:
                            try: _qtr_browser.close()
                            except Exception: pass
                        if _qtr_pw:
                            try: _qtr_pw.stop()
                            except Exception: pass

                if m1 and m2 and m3:
                    logger.set('Building PPT deck ...')
                    import pandas as pd
                    lp = results.get('labels_pivot')
                    dp = results.get('defects_pivot')
                    xp = results.get('exec_excel')
                    missing = []
                    if not lp:          missing.append('labels_pivot=None')
                    elif not lp.exists(): missing.append('labels_pivot NOT ON DISK: %s' % lp)
                    if not dp:          missing.append('defects_pivot=None')
                    elif not dp.exists(): missing.append('defects_pivot NOT ON DISK: %s' % dp)
                    if not xp:          missing.append('exec_excel=None')
                    elif not xp.exists(): missing.append('exec_excel NOT ON DISK: %s' % xp)
                    if missing:
                        for m_item in missing:
                            print('[WARN] PPT data missing: %s — will skip' % m_item)
                    labels_df = pd.read_excel(lp) if lp and lp.exists() else pd.DataFrame()
                    defects_df = pd.read_excel(dp) if dp and dp.exists() else pd.DataFrame()
                    qtr_df = pd.read_excel(xp) if xp and xp.exists() else pd.DataFrame()
                    release_value = 'INTEGRATION_PROGRESSION_%s' % iteration
                    def ensure_release_first(df):
                        df = df.copy()
                        if 'Release' not in df.columns:
                            df.insert(0, 'Release', release_value)
                        else:
                            cols = ['Release'] + [c for c in df.columns if c != 'Release']
                            df = df[cols]
                        return df
                    labels_df = ensure_release_first(labels_df)
                    defects_df = ensure_release_first(defects_df)
                    qtr_df = ensure_release_first(qtr_df)
                    prs = Presentation()
                    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
                    policy = LayoutPolicy(rows_per_part=10, top_in=1.05, side_in=0.90, footer_in=2.00, body_pt=8, header_pt=10)
                    logo = None
                    add_title_slide(prs, logo)
                    # --- Slide 2: Delivery Updates from template ---
                    try:
                        add_slide2_delivery_updates(prs, iteration, labels_df=labels_df, defects_df=defects_df, exec_df=qtr_df, insert_index=1)
                    except Exception as e2:
                        print('[WARN] Slide 2 insertion failed: %s' % e2)
                    # Add Total column + Total row to defects_df for Slide 3
                    defects_df_totals = defects_df.copy()
                    num_cols = [c for c in defects_df_totals.columns if c not in ('Release', 'Status')]
                    if num_cols:
                        defects_df_totals[num_cols] = defects_df_totals[num_cols].apply(pd.to_numeric, errors='coerce').fillna(0).astype(int)
                        defects_df_totals['Total'] = defects_df_totals[num_cols].sum(axis=1)
                        totals_row = {c: '' for c in defects_df_totals.columns}
                        totals_row['Status'] = 'Total'
                        for c in num_cols + ['Total']:
                            totals_row[c] = int(defects_df_totals[c].sum())
                        defects_df_totals = pd.concat([defects_df_totals, pd.DataFrame([totals_row])], ignore_index=True)
                    if not defects_df.empty:
                        add_table_parts(prs, policy, defects_df_totals, 'Quality Improvement: Progression Defect Summary', logo)
                    if not labels_df.empty:
                        add_table_parts(prs, policy, labels_df, 'Quality Improvement: Test Case Review Summary', logo)
                    if not qtr_df.empty:
                        add_table_parts(prs, policy, qtr_df, 'QMetry Test Execution Summary', logo)
                    # V3.0: Add cycle report slides before Thank You
                    cycle_data = results.get('cycle_reports')
                    if cycle_data:
                        try:
                            add_cycle_report_slides(prs, policy, cycle_data, logo)
                            print('[OK] Cycle report slides added')
                        except Exception as cre:
                            print('[WARN] Cycle slides failed: %s' % cre)
                    # V3.0: Features in Scope table
                    lbl_path = results.get('label_excel')
                    if lbl_path and lbl_path.exists():
                        try:
                            features_df = pd.read_excel(lbl_path)
                            add_table_parts_colored(prs, policy, features_df, 'Features in Scope', logo)
                            print('[OK] Features in Scope slides added')
                        except Exception as fe:
                            print('[WARN] Features slides failed: %s' % fe)
                    add_thank_you(prs, logo)
                    out_ppt = ARTIFACTS / ('MDA_QA_TMobile_INTG_TOSCA_Weekly_Status_%s.pptx' % datetime.now().strftime('%Y%m%d_%H%M%S'))
                    prs.save(out_ppt)
                    results['ppt'] = out_ppt
                    logger.set('✅ Deck built')
                    print('[DONE] Deck built => %s' % out_ppt)

                _elapsed = time.time() - _run_t0
                _m, _s = divmod(int(_elapsed), 60)
                print('\n⏱ Total execution time: %dm %ds' % (_m, _s))
                logger.set('✅ Completed — %dm %ds' % (_m, _s))
                run_success = True
            except Exception as e:
                _elapsed = time.time() - _run_t0
                _m, _s = divmod(int(_elapsed), 60)
                print('\n⏱ Total execution time: %dm %ds (failed)' % (_m, _s))
                logger.set('❌ Error — %dm %ds' % (_m, _s))
                traceback.print_exc()
                st.error('Error: %s' % e)

        if run_success:
            completed = []
            if results.get('labels_pivot'): completed.append('QMetry Labels Pivot ✅')
            if results.get('defects_pivot'): completed.append('Jira Defects Pivot ✅')
            if results.get('exec_excel'): completed.append('QMetry Test Execution Excel ✅')
            if results.get('cycle_reports'): completed.append('QMetry Cycle Pie Charts ✅')
            if results.get('ppt'): completed.append('Weekly Status Deck (PPTX) ✅')
            ss['last_summary'] = {
                'iteration': iteration,
                'ts': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'completed': completed,
                'labels_pivot': str(results.get('labels_pivot') or ''),
                'defects_pivot': str(results.get('defects_pivot') or ''),
                'exec_excel': str(results.get('exec_excel') or ''),
                'ppt': str(results.get('ppt') or ''),
            }
            ss['reset_mods'] = True
            st.toast('Module selections cleared')
            st.rerun()

# ---------------- Outputs section (same-run) ----------------
if results:
    with outputs_container:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.subheader('Outputs')
        oc1, oc2 = st.columns(2)
        with oc1:
            if results.get('labels_pivot') and results['labels_pivot'].exists():
                p = results['labels_pivot']
                st.write('QMetry Labels Pivot (Excel)')
                st.download_button('Download Labels Pivot', data=p.read_bytes(), file_name=p.name, key='dl_labels')
            if results.get('defects_pivot') and results['defects_pivot'].exists():
                p = results['defects_pivot']
                st.write('Jira Defects Pivot (Excel)')
                st.download_button('Download Defects Pivot', data=p.read_bytes(), file_name=p.name, key='dl_defects')
        with oc2:
            if results.get('exec_excel') and results['exec_excel'].exists():
                p = results['exec_excel']
                st.write('QMetry Test Execution Summary (Excel)')
                st.download_button('Download Test Report Excel', data=p.read_bytes(), file_name=p.name, key='dl_exec')
            if results.get('ppt') and results['ppt'].exists():
                p = results['ppt']
                st.write('Weekly Status Deck (PPTX)')
                st.download_button('Download PPT', data=p.read_bytes(), file_name=p.name, key='dl_ppt')
        st.markdown("</div>", unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True)  # end content-wrap

st.markdown("<div style='text-align:center;padding:12px 0;color:#b8c0cc;font-size:11px;font-weight:500;letter-spacing:0.5px;'>"
            "MDA Jira Dashboard &bull; V3.0 &bull; Electric Teal</div>", unsafe_allow_html=True)
